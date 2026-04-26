"""Generate Medical Imaging Coursework presentation with speaker notes."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage
import os

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
BASE = "/Users/basiakoch/DIS/medical/coursework/inverse-vision"
F    = BASE + "/notebooks/report_figures"
F11  = F + "/exercise_1_1"
F12  = F + "/exercise_1_2"
F13  = F + "/exercise_1_3"
F21  = F + "/exercise_2_1"
F22  = F + "/exercise_2_2"

# ---------------------------------------------------------------------------
# PRESENTATION SETUP
# ---------------------------------------------------------------------------
prs = Presentation()
W = Inches(13.33)
H = Inches(7.5)
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ---------------------------------------------------------------------------
# COLOUR PALETTE
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x0d, 0x1b, 0x2a)
NAVY2     = RGBColor(0x13, 0x2a, 0x47)
WHITE     = RGBColor(0xff, 0xff, 0xff)
CT        = RGBColor(0x15, 0x65, 0xC0)
CT_L      = RGBColor(0xe3, 0xf2, 0xfd)
MRI       = RGBColor(0x2E, 0x7D, 0x32)
MRI_L     = RGBColor(0xe8, 0xf5, 0xe9)
SEG       = RGBColor(0xbf, 0x36, 0x00)
SEG_L     = RGBColor(0xfb, 0xe9, 0xe7)
GRAY      = RGBColor(0x44, 0x55, 0x66)
LGRAY     = RGBColor(0xf0, 0xf4, 0xf8)
CT_BADGE  = RGBColor(0xbb, 0xde, 0xff)
MRI_BADGE = RGBColor(0xb8, 0xf0, 0xc4)
SEG_BADGE = RGBColor(0xff, 0xcc, 0xaa)

# ---------------------------------------------------------------------------
# LOW-LEVEL HELPERS
# ---------------------------------------------------------------------------
def new_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, color, border=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = color
    else:
        s.line.fill.background()
    return s

def txt(slide, content, l, t, w, h,
        size=14, color=NAVY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    lines = content.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size      = Pt(size)
        run.font.color.rgb = color
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.name      = "Calibri"
    return tb

def pic(slide, path, l, t, max_w, max_h):
    """Insert image preserving aspect ratio, centred within the target box."""
    if not path or not os.path.exists(path):
        return None
    with PILImage.open(path) as im:
        img_w_px, img_h_px = im.size
    img_aspect = img_w_px / img_h_px
    box_aspect  = max_w   / max_h
    if img_aspect > box_aspect:
        new_w = max_w
        new_h = int(max_w / img_aspect)
    else:
        new_h = max_h
        new_w = int(max_h * img_aspect)
    off_l = (max_w - new_w) // 2
    off_t = (max_h - new_h) // 2
    return slide.shapes.add_picture(path, l + off_l, t + off_t, new_w, new_h)

def notes(slide, text):
    """Add speaker notes to a slide (only visible in Presenter View)."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    paras = text.strip().split('\n')
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = para
        run.font.size = Pt(12)
        run.font.name = "Calibri"

# ---------------------------------------------------------------------------
# TABLE HELPER  — APA style built from shapes (guaranteed rendering)
# ---------------------------------------------------------------------------
RULE_H  = int(Pt(2.0))   # thick rule thickness in EMU
THIN_H  = int(Pt(0.9))   # thin rule thickness in EMU
ROW_H   = Inches(0.42)   # data row height
HDR_H   = Inches(0.48)   # header row height

def make_table(slide, headers, rows, col_fracs, l, t, w,
               accent=None, gd_rows=None):
    """
    APA-style table: thick top rule, bold header, thin rule under header,
    data rows, thick bottom rule. No vertical lines. Built from shapes.
    accent  — colour used for highlighted rows and header text label
    gd_rows — set of 0-based data-row indices to highlight (light green bg)
    Returns the y-coordinate of the bottom edge.
    """
    if accent is None:
        accent = CT

    col_widths = [int(w * f) for f in col_fracs]

    # ── thick top rule ──
    box(slide, l, t, w, RULE_H, NAVY)
    y = t + RULE_H

    # ── header row ──
    x = l
    for j, hdr in enumerate(headers):
        cw = col_widths[j]
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        txt(slide, hdr,
            x + Inches(0.08), y + Inches(0.04),
            cw - Inches(0.10), HDR_H - Inches(0.06),
            size=12, color=NAVY, bold=True, align=align)
        x += cw
    y += HDR_H

    # ── thin rule under header ──
    box(slide, l, y, w, THIN_H, NAVY)
    y += THIN_H

    # ── data rows ──
    for i, row in enumerate(rows):
        highlight = bool(gd_rows and i in gd_rows)
        row_bg = RGBColor(0xe8, 0xf5, 0xe9) if highlight else (
                 RGBColor(0xf7, 0xf7, 0xf7) if i % 2 == 0 else WHITE)
        # background strip
        box(slide, l, y, w, ROW_H, row_bg)
        x = l
        for j, val in enumerate(row):
            cw = col_widths[j]
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            color = accent if highlight else (NAVY if j == 0 else GRAY)
            txt(slide, val,
                x + Inches(0.08), y + Inches(0.04),
                cw - Inches(0.10), ROW_H - Inches(0.06),
                size=11, color=color, bold=highlight, align=align)
            x += cw
        y += ROW_H

    # ── thick bottom rule ──
    box(slide, l, y, w, RULE_H, NAVY)
    return y + RULE_H   # bottom edge

# ---------------------------------------------------------------------------
# SLIDE FACTORIES
# ---------------------------------------------------------------------------
def content_slide(accent, title, sub=None):
    slide = new_slide()
    set_bg(slide, WHITE)
    box(slide, 0, 0, W, Inches(0.10), accent)
    box(slide, 0, H - Inches(0.05), W, Inches(0.05), accent)
    txt(slide, title,
        Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.68),
        size=27, color=accent, bold=True)
    if sub:
        txt(slide, sub,
            Inches(0.35), Inches(0.78), Inches(12.6), Inches(0.38),
            size=13, color=GRAY, italic=True)
    return slide

def section_slide(bg_color, badge_color, module_label, heading, sub):
    slide = new_slide()
    set_bg(slide, bg_color)
    txt(slide, module_label,
        Inches(0.5), Inches(2.05), Inches(12), Inches(0.55),
        size=18, color=badge_color, align=PP_ALIGN.CENTER)
    txt(slide, heading,
        Inches(0.5), Inches(2.65), Inches(12), Inches(1.5),
        size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, sub,
        Inches(0.5), Inches(4.35), Inches(12), Inches(0.6),
        size=17, color=badge_color, align=PP_ALIGN.CENTER)
    return slide

# ===========================================================================
# S L I D E S
# ===========================================================================

# ── 1. TITLE ────────────────────────────────────────────────────────────────
s = new_slide()
set_bg(s, NAVY)
box(s, 0, Inches(2.35), W, Inches(2.95), NAVY2)
box(s, 0, Inches(2.35), W, Inches(0.07), CT)
box(s, 0, Inches(5.30), W, Inches(0.07), CT)
txt(s, "Medical Imaging Coursework",
    Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
    size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Inverse Problems in CT Reconstruction, MRI Denoising & Segmentation",
    Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.65),
    size=19, color=CT_BADGE, align=PP_ALIGN.CENTER)
txt(s, "Barbara Koch   ·   MPhil Data-Intensive Science   ·   University of Cambridge   ·   Lent 2026",
    Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.5),
    size=14, color=RGBColor(0x78, 0x90, 0xa8), align=PP_ALIGN.CENTER)
box(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.03), RGBColor(0x30, 0x60, 0x90))
notes(s, """Welcome. This presentation covers my Medical Imaging Coursework from Lent 2026.
I'll take you through three modules: CT reconstruction under dose and angular constraints, MRI multi-coil reconstruction and denoising, and a critical literature review on medical image segmentation.
The thread connecting all three is inverse problems — recovering hidden structure from indirect, noisy measurements.""")

# ── 2. PROBLEM OVERVIEW ──────────────────────────────────────────────────────
s = content_slide(NAVY, "The Problem: Inverse Problems in Medical Imaging")
cols = [
    (CT,  "Module 1", "CT Reconstruction",
     "X-ray projections measured at many angles.\nGoal: recover attenuation image from noisy sinogram.",
     "y = Ax + noise"),
    (MRI, "Module 2", "MRI Denoising",
     "Complex k-space data acquired from 6 receiver coils.\nGoal: reconstruct clean magnitude image.",
     "k = F[ρ] + noise"),
    (SEG, "Module 3", "Segmentation",
     "Assign a structure label to each pixel.\nDeep learning (R2U-Net) vs atlas-based methods.",
     "image → mask"),
]
for i, (color, mod, title, desc, formula) in enumerate(cols):
    x = Inches(0.30 + i * 4.34)
    y = Inches(1.15)
    w = Inches(4.12)
    h = Inches(5.65)
    box(s, x, y, w, h, LGRAY)
    box(s, x, y, w, Inches(0.10), color)
    txt(s, mod,   x+Inches(0.12), y+Inches(0.18), w-Inches(0.25), Inches(0.38), size=12, color=color, bold=True)
    txt(s, title, x+Inches(0.12), y+Inches(0.55), w-Inches(0.25), Inches(0.55), size=18, color=NAVY, bold=True)
    txt(s, desc,  x+Inches(0.12), y+Inches(1.15), w-Inches(0.25), Inches(2.4),  size=13, color=GRAY)
    box(s, x+Inches(0.12), y+Inches(3.85), w-Inches(0.25), Inches(0.6), color)
    txt(s, formula, x+Inches(0.12), y+Inches(3.88), w-Inches(0.25), Inches(0.55),
        size=15, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
txt(s, "All three are ill-posed inverse problems — many candidate solutions fit the data",
    Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.38),
    size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, """All three modules share a common mathematical structure: we observe indirect noisy measurements and need to invert that process.
In CT: the Radon transform A maps the image x to a sinogram; we observe a noisy version and want to recover x.
In MRI: the Fourier transform F maps tissue magnetisation to k-space; we observe noisy k-space and invert with IFFT.
In segmentation: we want to map from image pixels to structure labels — here the inverse is learned or atlas-based.
All three are ill-posed: the data alone doesn't uniquely determine the answer, so we need additional assumptions (regularisation, noise models, atlases).""")

# ── 3. CT SECTION HEADER ────────────────────────────────────────────────────
s = section_slide(CT, CT_BADGE, "MODULE  1",
                  "CT Reconstruction",
                  "Dose reduction  ·  Limited angle  ·  Filters  ·  Iterative methods")
notes(s, """Starting with Module 1: CT Reconstruction.
This module has three exercises: 1.1 — dose reduction study comparing FBP vs Gradient Descent; 1.2 — limited-angle tomography; 1.3 — FBP filter comparison and subset GD analysis.
The overarching question: how do different acquisition constraints affect reconstruction quality, and which algorithm handles them best?""")

# ── 4. CT PROBLEM SETUP ──────────────────────────────────────────────────────
s = content_slide(CT, "CT: The Forward Model",
                  "X-ray projections through tissue  →  sinogram  →  reconstruct attenuation image")
bullets = [
    "Phantom: 512×512 chest CT, masked to circular FOV",
    "Parallel-beam CT: angles sampled uniformly over [0°, 180°]",
    "Beer–Lambert law:  I = I₀ exp(−Ax)",
    "Gaussian detector noise σ = 0.05 + Poisson photon noise",
    "Noisy log-sinogram: b = −log(I_noisy / I₀) · p_max",
    "Inverse problem:  min ‖Ax − b‖²",
]
for i, b in enumerate(bullets):
    txt(s, "•  " + b, Inches(0.35), Inches(1.25 + i * 0.83), Inches(6.1), Inches(0.78), size=14, color=GRAY)
box(s, Inches(0.35), Inches(6.35), Inches(6.1), Inches(0.72), CT_L)
txt(s, "9 conditions: I₀ ∈ {10⁵, 10³, 10²}  ×  n ∈ {360, 90, 20}",
    Inches(0.45), Inches(6.40), Inches(5.9), Inches(0.62),
    size=16, color=CT, italic=True, align=PP_ALIGN.CENTER)
pic(s, F11 + "/figure3a_sinograms_gaussian_only.png", Inches(6.5), Inches(1.15), Inches(6.65), Inches(5.95))
txt(s, "Sinograms with Gaussian noise only  (all 9 conditions)  —  Fig. 1 in report",
    Inches(6.5), Inches(7.1), Inches(6.65), Inches(0.35),
    size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, """What I actually did:
— Loaded the 512×512 chest CT phantom and masked it to the inscribed circle so only the reliably reconstructable region is kept.
— Simulated parallel-beam CT with the Radon transform. Angles are uniformly spaced over [0°, 180°].
— Added noise in two steps: first Gaussian detector noise (σ=0.05) to model electronic noise, then Poisson sampling to model photon-counting statistics. This two-step model is more realistic than Poisson alone.
— Converted noisy transmission counts back to projection values using the log transform — that's the sinogram we try to invert.
— I₀ controls how many photons hit the detector: high I₀=10⁵ means low noise; low I₀=10² means heavy noise.
The image on the right shows the Gaussian-noise-only sinograms (Fig. 1 in the report) — an auxiliary diagnostic to isolate the detector noise term before adding Poisson photon noise. This gives the audience a first look at what sinograms are before the full noisy versions appear on the next slide.""")

# ── 4b. GAUSSIAN-NOISE-ONLY SINOGRAMS (Fig. 1) ──────────────────────────────
s = content_slide(CT, "Sinograms with Gaussian Noise Only  (Fig. 1)",
                  "Auxiliary diagnostic — isolates detector noise before Poisson photon noise is added")
pic(s, F11 + "/figure3a_sinograms_gaussian_only.png",
    Inches(0.25), Inches(1.15), Inches(12.83), Inches(5.65))
box(s, Inches(0.25), Inches(6.86), Inches(6.0), Inches(0.48), CT_L)
txt(s, "↓ decreasing dose  (I₀ = 10⁵ → 10³ → 10²)",
    Inches(0.35), Inches(6.88), Inches(5.8), Inches(0.44), size=12, color=CT, bold=True)
box(s, Inches(6.6), Inches(6.86), Inches(6.48), Inches(0.48), CT_L)
txt(s, "← fewer projection angles  (360 → 90 → 20 views)",
    Inches(6.7), Inches(6.88), Inches(6.28), Inches(0.44), size=12, color=CT, bold=True)
notes(s, """Figure 1 in the report — Gaussian-noise-only sinograms for all 9 conditions.
This is an auxiliary diagnostic: I generated this version by applying only the Gaussian detector noise step (σ=0.05) and skipping the Poisson photon sampling step.
Why show this? To isolate the effect of each noise source separately. Here you see that Gaussian noise adds a smooth, intensity-scaled graininess. Comparing this to Figure 2 (which adds Poisson on top) shows the extra, count-dependent fluctuations that Poisson noise introduces.
The structure is still clearly visible in all 9 conditions — Gaussian alone is relatively benign. The damaging noise at low dose is mainly the Poisson component, because fewer photons means each count fluctuation is a larger fraction of the signal.""")

# ── 5. SINOGRAM GRID ────────────────────────────────────────────────────────
s = content_slide(CT, "9 Noisy Sinograms: Dose × Angular Sampling",
                  "Rows: I₀ = 10⁵ → 10³ → 10²  |  Columns: 360 → 90 → 20 views")
pic(s, F11 + "/figure3_noisy_sinograms_clean.png",
    Inches(0.25), Inches(1.15), Inches(12.83), Inches(5.65))
box(s, Inches(0.25), Inches(6.86), Inches(6.0), Inches(0.48), CT_L)
txt(s, "↓ decreasing dose (more noise per ray)",
    Inches(0.35), Inches(6.88), Inches(5.8), Inches(0.44), size=12, color=CT, bold=True)
box(s, Inches(6.6), Inches(6.86), Inches(6.48), Inches(0.48), CT_L)
txt(s, "← fewer projection angles (undersampling artefacts)",
    Inches(6.7), Inches(6.88), Inches(6.28), Inches(0.44), size=12, color=CT, bold=True)
notes(s, """This is the 3×3 experimental grid I generated — Figure 2 in my report.
Each cell is the noisy sinogram for one (I₀, n_views) condition.
Down the rows: I₀ drops from 10⁵ → 10³ → 10², and you can see the sinogram becoming grainier — more photon noise corrupts each ray measurement.
Across the columns: the number of projections drops from 360 → 90 → 20. With fewer angles the sinogram has fewer rows, which means less angular information for reconstruction.
The two effects are independent: dose controls noise amplitude per ray, while angular count controls how densely the Radon domain is sampled. The worst case is bottom-right: very noisy AND very sparse.""")

# ── 6. RECONSTRUCTION METHODS ────────────────────────────────────────────────
s = content_slide(CT, "Two Reconstruction Methods")
box(s, Inches(0.30), Inches(1.15), Inches(6.0), Inches(5.65), CT_L)
box(s, Inches(0.30), Inches(1.15), Inches(6.0), Inches(0.10), CT)
txt(s, "Filtered Backprojection  (FBP)",
    Inches(0.45), Inches(1.30), Inches(5.7), Inches(0.52), size=19, color=CT, bold=True)
for i, b in enumerate(["Non-iterative — direct analytic formula",
                        "Ram-Lak ramp filter in Fourier domain",
                        "Backprojects filtered projections",
                        "Fast  (∼1 s per reconstruction)",
                        "Amplifies noise — degrades badly at low dose"]):
    txt(s, "•  " + b, Inches(0.45), Inches(1.92 + i * 0.68), Inches(5.7), Inches(0.62), size=13, color=GRAY)
box(s, Inches(0.45), Inches(5.35), Inches(5.7), Inches(0.7), CT)
txt(s, "x = Aᵀ (ramp ∗ b)",
    Inches(0.45), Inches(5.38), Inches(5.7), Inches(0.64),
    size=18, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
box(s, Inches(7.03), Inches(1.15), Inches(6.0), Inches(5.65), MRI_L)
box(s, Inches(7.03), Inches(1.15), Inches(6.0), Inches(0.10), MRI)
txt(s, "Gradient Descent  (GD)",
    Inches(7.18), Inches(1.30), Inches(5.7), Inches(0.52), size=19, color=MRI, bold=True)
for i, b in enumerate(["Iterative — 200 iterations",
                        "Minimises ‖Ax − b‖² directly",
                        "Initialised from zero, step size γ = 0.001",
                        "No regularisation / non-negativity constraint",
                        "Robust to noise — outperforms FBP in all 9 conditions"]):
    txt(s, "•  " + b, Inches(7.18), Inches(1.92 + i * 0.68), Inches(5.7), Inches(0.62), size=13, color=GRAY)
box(s, Inches(7.18), Inches(5.35), Inches(5.7), Inches(0.7), MRI)
txt(s, "xᵏ⁺¹ = xᵏ − 0.001 · Aᵀ(Axᵏ − b)",
    Inches(7.18), Inches(5.38), Inches(5.7), Inches(0.64),
    size=18, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
txt(s, "vs", Inches(6.17), Inches(3.55), Inches(0.99), Inches(0.55),
    size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, """FBP is the classical analytic approach. It works by filtering each 1D projection with the ramp filter in the frequency domain (to correct the backprojection blurring), then smearing the filtered projections back across the image. The ramp filter boosts high frequencies — which is great for sharpness but terrible for noise, because noise lives at high frequencies.
Gradient Descent is iterative. Each iteration I compute the residual Ax^k - b (how well the current estimate explains the data), backproject it through A^T, and subtract a small step in that direction. After 200 iterations, the estimate has converged toward the least-squares minimum.
The key difference: FBP makes a single pass through the data; GD makes 200 passes and adapts. This is why GD can partially 'see through' noise — it progressively refines the estimate rather than making one noisy guess.""")

# ── 6b. TABLE 1 — PSNR / SSIM ALL 9 CONDITIONS ──────────────────────────────
s = content_slide(CT, "Results: PSNR (dB) / SSIM — All 9 Conditions  (Table 1)",
                  "GD outperforms FBP in every single condition  ·  Bold = GD (better method)")
t1_hdrs = ["I₀", "Method", "PSNR\n360 views", "SSIM\n360 views",
           "PSNR\n90 views", "SSIM\n90 views", "PSNR\n20 views", "SSIM\n20 views"]
t1_rows = [
    ["10⁵", "FBP",   "36.91", "0.891", "29.69", "0.642", "20.13", "0.311"],
    ["10⁵", "GD ✓",  "37.08", "0.954", "35.15", "0.902", "28.14", "0.673"],
    ["10³", "FBP",   "17.30", "0.264", "11.27", "0.206", " 4.40", "0.185"],
    ["10³", "GD ✓",  "27.05", "0.517", "23.45", "0.377", "23.60", "0.369"],
    ["10²", "FBP",   " 7.20", "0.190", " 1.19", "0.182", "−5.67", "0.180"],
    ["10²", "GD ✓",  "17.33", "0.250", "13.62", "0.211", "15.15", "0.213"],
]
make_table(s, t1_hdrs, t1_rows,
           [0.09, 0.09, 0.135, 0.135, 0.135, 0.135, 0.135, 0.135],
           Inches(0.30), Inches(1.20), Inches(12.73),
           gd_rows={1, 3, 5}, accent=MRI)
box(s, Inches(0.30), Inches(5.55), Inches(12.73), Inches(0.60), CT_L)
txt(s, "GD dominates all 9 conditions without exception. "
       "The gap narrows at I₀=10⁵, 360 views (37.08 vs 36.91 dB) "
       "but widens dramatically at low dose — worst case: GD 17.33 dB vs FBP 7.20 dB at I₀=10², 360 views.",
    Inches(0.40), Inches(5.58), Inches(12.53), Inches(0.55),
    size=12, color=CT, align=PP_ALIGN.CENTER)
box(s, Inches(0.30), Inches(6.28), Inches(6.20), Inches(0.55), LGRAY)
txt(s, "FBP PSNR collapses to −5.67 dB at I₀=10², 20 views\n→ clinically unusable",
    Inches(0.40), Inches(6.30), Inches(6.00), Inches(0.51),
    size=12, color=RGBColor(0xcc, 0x00, 0x00), bold=True)
box(s, Inches(6.83), Inches(6.28), Inches(6.20), Inches(0.55), LGRAY)
txt(s, "GD at same condition: 15.15 dB — anatomy still partially visible",
    Inches(6.93), Inches(6.30), Inches(6.00), Inches(0.51),
    size=12, color=MRI, bold=True)
notes(s, """Table 1 from the report: PSNR (dB) and SSIM for all nine acquisition conditions.
GD wins in all 9 out of 9 conditions — this is the headline finding.
Key numbers to mention:
- At I₀=10⁵, 360 views: FBP 36.91 vs GD 37.08 — almost identical, both excellent.
- At I₀=10², 20 views: FBP −5.67 dB (worse than predicting a constant!), GD 15.15 dB — clinically usable.
- The gap is largest at low dose, high view count (I₀=10², 360 views): FBP 7.20 vs GD 17.33 dB — 10 dB difference.
The green-highlighted rows are GD. FBP degrades sharply when either dose or view count drops; GD degrades more gracefully because it iteratively refines the estimate rather than making a single noisy pass.""")

# ── 7. RESULTS HIGH DOSE ─────────────────────────────────────────────────────
s = content_slide(CT, "Results  —  High Dose  (I₀ = 10⁵)",
                  "FBP performs well  ·  GD slightly better across all view counts")
pic(s, F11 + "/recon_I0_1e5.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
box(s, Inches(0.20), Inches(7.05), Inches(5.0), Inches(0.38), CT_L)
txt(s, "FBP  —  PSNR 36.9 / 29.7 / 20.1 dB  (360 / 90 / 20 views)",
    Inches(0.30), Inches(7.06), Inches(4.9), Inches(0.36), size=11, color=CT)
box(s, Inches(7.0), Inches(7.05), Inches(5.5), Inches(0.38), MRI_L)
txt(s, "GD  —  PSNR 37.1 / 35.2 / 28.1 dB  (360 / 90 / 20 views)",
    Inches(7.1), Inches(7.06), Inches(5.3), Inches(0.36), size=11, color=MRI)
notes(s, """At high dose, both methods work well — this is the easy case.
At 360 views, FBP gives 36.9 dB and GD gives 37.1 dB — almost identical, the gap is under 0.2 dB.
But look at what happens as views drop. At 20 views, FBP drops to 20.1 dB while GD gives 28.1 dB — a gap of 8 dB. GD clearly handles angular undersampling better even when noise is low.
The columns of figures: FBP images (cols 1&2) and GD images (cols 3&4), with error maps next to each. GD error maps are much darker, especially at 20 views. The errors in GD are concentrated at edges; FBP errors are spread across the image as background noise.""")

# ── 8. RESULTS MEDIUM DOSE ───────────────────────────────────────────────────
s = content_slide(CT, "Results  —  Medium Dose  (I₀ = 10³)",
                  "FBP degrades noticeably  ·  GD maintains structure")
pic(s, F11 + "/recon_I0_1e3.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
box(s, Inches(0.20), Inches(7.05), Inches(5.0), Inches(0.38), CT_L)
txt(s, "FBP  —  PSNR 17.3 / 11.3 / 4.4 dB",
    Inches(0.30), Inches(7.06), Inches(4.9), Inches(0.36), size=11, color=CT)
box(s, Inches(7.0), Inches(7.05), Inches(5.5), Inches(0.38), MRI_L)
txt(s, "GD  —  PSNR 27.0 / 23.5 / 23.6 dB  —  far more robust",
    Inches(7.1), Inches(7.06), Inches(5.3), Inches(0.36), size=11, color=MRI)
notes(s, """At medium dose I₀=10³, the divergence between the methods becomes clinically significant.
FBP at 360 views: 17.3 dB — you can still see anatomy but there's visible noise. At 20 views: 4.4 dB — barely interpretable.
GD at 360 views: 27.0 dB — much cleaner. At 20 views: 23.6 dB — notice this is actually SIMILAR to 360 views. GD is remarkably stable across angular counts at this dose level.
That's the key insight for this slide: GD's robustness is especially visible in the view-count dimension. FBP is sensitive to both noise AND undersampling simultaneously; GD handles them more gracefully.""")

# ── 9. RESULTS LOW DOSE ──────────────────────────────────────────────────────
s = content_slide(CT, "Results  —  Low Dose  (I₀ = 10²)  —  FBP Fails",
                  "GD partially recovers anatomy  ·  FBP is noise-dominated and clinically unusable")
pic(s, F11 + "/recon_I0_1e2.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.7))
box(s, Inches(0.20), Inches(6.92), Inches(5.7), Inches(0.46), RGBColor(0xff, 0xdd, 0xdd))
txt(s, "FBP  PSNR  7.2 / 1.2 / −5.7 dB   —   unusable at low dose",
    Inches(0.30), Inches(6.94), Inches(5.5), Inches(0.42),
    size=12, color=RGBColor(0xcc, 0x00, 0x00), bold=True)
box(s, Inches(7.0), Inches(6.92), Inches(6.0), Inches(0.46), RGBColor(0xcc, 0xee, 0xcc))
txt(s, "GD  PSNR  17.3 / 13.6 / 15.1 dB   —   anatomy partially preserved",
    Inches(7.1), Inches(6.94), Inches(5.8), Inches(0.42),
    size=12, color=RGBColor(0x00, 0x66, 0x00), bold=True)
notes(s, """This is the most dramatic result. At I₀=10², the photon count is so low that noise completely dominates the sinogram.
FBP at 20 views: −5.7 dB PSNR — negative PSNR means the reconstruction is WORSE than just predicting a constant image. It is clinically worthless. Even at 360 views FBP only manages 7.2 dB.
GD at 20 views: 15.1 dB — you can still see the heart, the lungs, the major vessels. At 360 views: 17.3 dB. GD has partially compensated for both the heavy noise and the sparse angles.
You can see this directly in the figures: FBP images are swamped with salt-and-pepper noise, while GD images, though blurrier, retain clear anatomical boundaries. This is the main result of Exercise 1.1.""")

# ── 10. METRICS HEATMAP ──────────────────────────────────────────────────────
s = content_slide(CT, "GD Dominates All 9 Conditions",
                  "PSNR and SSIM heatmaps across all dose × view combinations")
pic(s, F11 + "/metrics_heatmap.png",
    Inches(0.40), Inches(1.15), Inches(12.53), Inches(5.75))
box(s, Inches(0.35), Inches(6.98), Inches(12.63), Inches(0.42), CT_L)
txt(s, "GD outperforms FBP in every single combination  ·  "
       "Gap narrows only at I₀=10⁵, 360 views (37.1 vs 36.9 dB)  ·  "
       "Largest gap: I₀=10², 360 views (17.3 vs 7.2 dB)",
    Inches(0.45), Inches(7.00), Inches(12.43), Inches(0.38),
    size=12, color=CT, bold=True, align=PP_ALIGN.CENTER)
notes(s, """This heatmap is the summary figure for Exercise 1.1 — Figure 6 in the report.
Four heatmaps: FBP-PSNR, GD-PSNR, FBP-SSIM, GD-SSIM. Rows are I₀ (10⁵ top to 10² bottom), columns are views (360, 90, 20 left to right).
Green = good, red = bad. You can see immediately that GD's heatmaps are much greener, especially at low dose.
The single most important number: GD wins in all 9 out of 9 conditions without exception. The margin is tiny at the best case (37.1 vs 36.9 dB) but massive at the worst (17.3 vs 7.2 dB — a 10 dB gap).
This is Table 1 in my report quantified visually.""")

# ── 11. CLINICAL INTERPRETATION ──────────────────────────────────────────────
s = content_slide(CT, "Clinical Implications")
rows = [
    ("High dose + many views",
     "Best diagnostic quality  —  but highest radiation exposure to the patient"),
    ("Low dose / sparse views",
     "GD can partially recover anatomy  —  FBP completely fails below I₀=10³"),
    ("Dose–quality trade-off",
     "Iterative methods enable dose reduction without sacrificing diagnostic value"),
    ("Computational cost",
     "GD: slower (~200 iter), but clinically justified in low-dose protocols"),
]
for i, (title, desc) in enumerate(rows):
    y = Inches(1.25 + i * 1.43)
    box(s, Inches(0.30), y, Inches(12.73), Inches(1.30), LGRAY)
    box(s, Inches(0.30), y, Inches(0.09), Inches(1.30), CT)
    txt(s, title, Inches(0.50), y + Inches(0.10), Inches(12.4), Inches(0.45), size=17, color=CT, bold=True)
    txt(s, desc,  Inches(0.50), y + Inches(0.55), Inches(12.4), Inches(0.65), size=14, color=GRAY)
notes(s, """Connecting the results to the clinical context.
The core tension in clinical CT is between radiation dose and image quality. More photons = better images, but radiation is harmful — especially for paediatric patients or people needing repeated scans.
Low-dose CT protocols are increasingly standard, but they produce noisier sinograms. Our results show that GD can still reconstruct diagnostically useful images at I₀=10² where FBP completely fails. This means iterative methods could enable dose reduction of perhaps a factor of 10-100 without losing diagnostic value.
The downside is speed: GD needs 200 iterations per reconstruction vs a single pass for FBP. In practice, GPU acceleration and regularised iterative methods (like TV-regularised GD or ADMM) are used clinically precisely because of this quality advantage.""")

# ── 12. IMPROVEMENTS ─────────────────────────────────────────────────────────
s = content_slide(CT, "How to Improve CT Reconstruction Further")
items = [
    ("TV Regularisation",
     "min ‖Ax − b‖² + λ‖∇x‖₁",
     "Suppress streak artefacts\nPreserve sharp edges\nEliminate semi-convergence in GD"),
    ("Sinogram Pre-filtering",
     "Denoise b before inversion",
     "Low-pass or learned denoiser\nReduces noise amplified by ramp filter\nImproves FBP at low dose"),
    ("Deep Learning (post-processing)",
     "FBP(b) → CNN → clean image",
     "Learn domain-specific noise patterns\nU-Net correction of FBP artefacts\nFast at inference time"),
]
for i, (title, formula, desc) in enumerate(items):
    x = Inches(0.30 + i * 4.34)
    w = Inches(4.12)
    box(s, x, Inches(1.15), w, Inches(5.85), LGRAY)
    box(s, x, Inches(1.15), w, Inches(0.10), CT)
    txt(s, title,   x+Inches(0.12), Inches(1.30), w-Inches(0.25), Inches(0.55), size=17, color=CT, bold=True)
    box(s, x+Inches(0.12), Inches(1.90), w-Inches(0.25), Inches(0.60), CT_L)
    txt(s, formula, x+Inches(0.18), Inches(1.93), w-Inches(0.37), Inches(0.55),
        size=14, color=CT, italic=True, align=PP_ALIGN.CENTER)
    txt(s, desc,    x+Inches(0.12), Inches(2.60), w-Inches(0.25), Inches(3.8), size=13, color=GRAY)
notes(s, """Three directions I discussed in my report for improving upon unregularised GD.
TV regularisation: adds a total-variation penalty λ‖∇x‖₁ to the objective. This penalises large image gradients, so the algorithm is pushed toward piecewise-smooth solutions. It suppresses noise in flat regions while preserving sharp edges — ideal for CT where you want to resolve tissue boundaries. It also fixes the semi-convergence problem I observed: unregularised GD peaks around epoch 33 then slowly degrades as noise gets fitted; TV prevents this by acting as a stopping criterion embedded in the objective.
Sinogram pre-filtering: apply a denoiser (simple Gaussian or learned neural network) to the noisy sinogram b before inversion. This reduces the noise that FBP's ramp filter amplifies. The cost is some loss of spatial resolution.
Deep learning post-processing: run FBP first (fast), then apply a U-Net to clean up the result. This is what companies like GE and Siemens use in current clinical scanners.""")

# ── 13. LIMITED ANGLE — SETUP ────────────────────────────────────────────────
s = content_slide(CT, "Module 1.2: Limited-Angle Tomography",
                  "Angular range restricted: 180° → 120° → 40°  (360 projections each)")
pic(s, F12 + "/noisy_sinograms.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.70))
for i, lbl in enumerate(["180°  (full)", "120°", "40°  (severely limited)"]):
    txt(s, lbl, Inches(0.30 + i * 4.34), Inches(6.93), Inches(4.14), Inches(0.42),
        size=13, color=CT, bold=True, align=PP_ALIGN.CENTER)
notes(s, """Exercise 1.2 changes the experiment: instead of reducing the number of projections, I restrict the angular range.
In Exercise 1.1, I had n ∈ {360, 90, 20} projections spanning all 180°. Here I have 360 projections in each case, but the angular range shrinks: 180°, 120°, 40°.
These sinograms all have the same width (360 columns = 360 angular samples), but the horizontal axis now represents a smaller arc. The rows are I₀ as before.
The physical motivation: in some clinical situations you cannot rotate the X-ray source all the way around — e.g. interventional radiology where equipment is in the way, or C-arm CT systems that can't complete a full arc.""")

# ── 14. LIMITED ANGLE — RESULTS ──────────────────────────────────────────────
s = content_slide(CT, "Limited-Angle Reconstructions  (I₀ = 10⁵)",
                  "Directional elongation artefacts grow as the angular range narrows")
pic(s, F12 + "/recon_I0_1e5.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.70))
for i, (lbl, clr) in enumerate([
    ("180°  PSNR 37.1 dB — clean", RGBColor(0x00, 0x66, 0x00)),
    ("120°  PSNR 25.3 dB — artefacts", RGBColor(0xaa, 0x55, 0x00)),
    ("40°  PSNR 20.2 dB — elongated", RGBColor(0xcc, 0x00, 0x00)),
]):
    box(s, Inches(0.30 + i * 4.34), Inches(6.88), Inches(4.14), Inches(0.48), LGRAY)
    txt(s, lbl, Inches(0.35 + i * 4.34), Inches(6.90), Inches(4.05), Inches(0.44),
        size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
notes(s, """This is at high dose so noise is minimal — the degradation we see is entirely due to the limited angular range.
At 180° (full range): clean reconstruction, 37.1 dB GD.
At 120°: clear directional streaks appear. Structures that should be smooth acquire directional blur or smearing. PSNR drops to 25.3 dB.
At 40°: the image is severely elongated in one direction. Structures that should be approximately round (like the heart) appear stretched. PSNR 20.2 dB.
Crucially, this is qualitatively DIFFERENT from the noise artefacts in Exercise 1.1. There, reducing views caused noise-like degradation; here, restricting the angular range causes coherent directional distortion. The anatomy is still recognisable but geometrically wrong.""")

# ── 14b. LIMITED ANGLE — RESULTS I₀=10³ ─────────────────────────────────────
s = content_slide(CT, "Limited-Angle Reconstructions  (I₀ = 10³)",
                  "Noise compounds the missing-wedge artefacts at medium dose")
pic(s, F12 + "/recon_I0_1e3.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.70))
for i, (lbl, clr) in enumerate([
    ("180°  PSNR 27.1 dB", RGBColor(0x00, 0x66, 0x00)),
    ("120°  PSNR 23.5 dB — streaks + noise", RGBColor(0xaa, 0x55, 0x00)),
    ("40°  PSNR 19.9 dB — elongated + noisy", RGBColor(0xcc, 0x00, 0x00)),
]):
    box(s, Inches(0.30 + i * 4.34), Inches(6.88), Inches(4.14), Inches(0.48), LGRAY)
    txt(s, lbl, Inches(0.35 + i * 4.34), Inches(6.90), Inches(4.05), Inches(0.44),
        size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
notes(s, """At medium dose I₀=10³, both noise and missing-wedge artefacts are present simultaneously.
The 180° image is 27.1 dB — same as Exercise 1.1 at this dose with 360 views, confirming consistency.
At 120°: directional streaks from the missing wedge combine with noise, producing a messier degradation pattern than either effect alone.
At 40°: both elongation artefacts and noise are severe — the two degradation modes compound each other.
Key observation: compare rows in Table 2. GD still outperforms FBP at every angle, but the gap narrows at low dose because noise limits both methods.""")

# ── 14c. LIMITED ANGLE — RESULTS I₀=10² ─────────────────────────────────────
s = content_slide(CT, "Limited-Angle Reconstructions  (I₀ = 10²)  —  Noise Dominates",
                  "At low dose, noise overwhelms angular-range differences — metrics flatten")
pic(s, F12 + "/recon_I0_1e2.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.70))
for i, (lbl, clr) in enumerate([
    ("180°  PSNR 17.3 dB", RGBColor(0xaa, 0x55, 0x00)),
    ("120°  PSNR 17.2 dB — noise dominates", RGBColor(0xaa, 0x55, 0x00)),
    ("40°  PSNR 17.8 dB — metrics flat", RGBColor(0xaa, 0x55, 0x00)),
]):
    box(s, Inches(0.30 + i * 4.34), Inches(6.88), Inches(4.14), Inches(0.48), LGRAY)
    txt(s, lbl, Inches(0.35 + i * 4.34), Inches(6.90), Inches(4.05), Inches(0.44),
        size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
notes(s, """The striking result at low dose: GD metrics are almost flat across angular ranges — 17.3, 17.2, 17.8 dB.
Compare to high dose (I₀=10⁵): 37.1 → 25.3 → 20.2 dB as angle shrinks. At high dose angular range strongly controls quality.
At I₀=10², noise is so dominant that the missing-wedge effect is a secondary concern — noise swamps everything regardless of angular coverage.
The practical implication: if you are in a low-dose regime, expanding the angular range does not help much. The priority must be noise reduction (more photons, better detector, iterative denoising) rather than acquiring over a wider arc.""")

# ── 14d. TABLE 2 — LIMITED-ANGLE PSNR / SSIM ────────────────────────────────
s = content_slide(CT, "Results: PSNR (dB) / SSIM — Limited-Angle  (Table 2)",
                  "360 projections per arc  ·  At high dose angular range dominates  ·  At low dose noise dominates")
t2_hdrs = ["I₀", "Method", "PSNR\n180°", "SSIM\n180°",
           "PSNR\n120°", "SSIM\n120°", "PSNR\n40°", "SSIM\n40°"]
t2_rows = [
    ["10⁵", "FBP",   "36.91", "0.891", "22.21", "0.537", "14.59", "0.295"],
    ["10⁵", "GD ✓",  "37.08", "0.954", "25.30", "0.710", "20.23", "0.565"],
    ["10³", "FBP",   "17.30", "0.264", "16.09", "0.237", "13.05", "0.201"],
    ["10³", "GD ✓",  "27.05", "0.517", "23.46", "0.447", "19.92", "0.433"],
    ["10²", "FBP",   " 7.20", "0.190", " 7.06", "0.189", " 7.28", "0.184"],
    ["10²", "GD ✓",  "17.33", "0.250", "17.19", "0.245", "17.78", "0.254"],
]
make_table(s, t2_hdrs, t2_rows,
           [0.09, 0.09, 0.135, 0.135, 0.135, 0.135, 0.135, 0.135],
           Inches(0.30), Inches(1.20), Inches(12.73),
           gd_rows={1, 3, 5}, accent=MRI)
box(s, Inches(0.30), Inches(5.55), Inches(6.15), Inches(0.80), CT_L)
txt(s, "High dose (I₀=10⁵): PSNR drops 37→25→20 dB as arc shrinks\n→ angular range is the dominant factor",
    Inches(0.40), Inches(5.58), Inches(5.95), Inches(0.75), size=12, color=CT, bold=True)
box(s, Inches(6.88), Inches(5.55), Inches(6.15), Inches(0.80), LGRAY)
txt(s, "Low dose (I₀=10²): GD metrics ~flat (17.3 / 17.2 / 17.8 dB)\n→ noise dominates, angular range barely matters",
    Inches(6.98), Inches(5.58), Inches(5.95), Inches(0.75), size=12, color=GRAY, bold=True)
notes(s, """Table 2 from the report: PSNR/SSIM for limited-angle reconstructions (360 projections per arc).
The 180° row matches Exercise 1.1 exactly (same conditions, same noise seed) — confirming consistency.
Key finding at high dose (I₀=10⁵): PSNR drops sharply as arc shrinks — 37.1 → 25.3 → 20.2 dB for GD. The angular range is the dominant degradation factor here.
Key finding at low dose (I₀=10²): GD metrics are nearly flat across all arc lengths — 17.33, 17.19, 17.78 dB. Noise is so dominant that the missing-wedge effect becomes secondary.
This tells us: in a low-dose protocol, acquiring a wider angular range gives you almost no benefit. The bottleneck is photon statistics, not angular coverage.""")

# ── 14e. DISCUSSION: DIFFERENCES IN THE LIMITED-ANGLE CASE (Ex 1.2b) ─────────
s = content_slide(CT, "Discussion: Differences in the Limited-Angle Case  (Ex 1.2b)",
                  "Two ways to reduce dose — different degradation modes, different clinical impact")
rows_diff = [
    ("Effect",       "Exercise 1.1 — Reducing views\n(360 → 90 → 20 projections, full 180° arc)",
                     "Exercise 1.2 — Restricting angle\n(180° → 120° → 40°, 360 projections each)"),
    ("Artefact type","Noise-like: increased graininess,\nrandom pixel-level fluctuations",
                     "Directional: elongation, streaking\nalong unsampled orientations"),
    ("Anatomy",      "Still recognisable even at 20 views;\noverall structure preserved",
                     "Geometrically distorted — heart and\nbones appear stretched or smeared"),
    ("Root cause",   "Fewer measurements → higher\nnoise per reconstructed pixel",
                     "Missing wedge in Fourier space →\ncertain edge orientations unrecoverable"),
    ("GD advantage", "GD robust to noise; outperforms\nFBP especially at sparse views",
                     "GD also fails — cannot recover\nmissing frequency information"),
    ("Clinical note","Low view count: faster scan,\nlower dose, but noisy images",
                     "Limited arc: relevant in C-arm CT,\nDBS tomosynthesis, interventional imaging"),
]
col_fracs = [0.175, 0.4125, 0.4125]
headers_diff = ["", "Ex 1.1 — Dose / Views", "Ex 1.2 — Limited Angle"]
for i, (label, v11, v12) in enumerate(rows_diff):
    y = Inches(1.18 + i * 0.98)
    bg = LGRAY if i % 2 == 0 else WHITE
    box(s, Inches(0.30), y, Inches(12.73), Inches(0.93), bg)
    box(s, Inches(0.30), y, Inches(0.08), Inches(0.93), CT)
    txt(s, label, Inches(0.44), y + Inches(0.06), Inches(1.85), Inches(0.80),
        size=11, color=CT, bold=True)
    txt(s, v11,   Inches(2.60), y + Inches(0.06), Inches(4.95), Inches(0.80),
        size=11, color=GRAY)
    txt(s, v12,   Inches(7.65), y + Inches(0.06), Inches(5.20), Inches(0.80),
        size=11, color=GRAY)
# column headers
for j, (hdr, x, w2) in enumerate([
    ("", Inches(0.30), Inches(2.20)),
    ("Ex 1.1 — Dose / Views", Inches(2.55), Inches(5.10)),
    ("Ex 1.2 — Limited Angle", Inches(7.60), Inches(5.40)),
]):
    box(s, x, Inches(1.05), w2, Inches(0.10), CT if j > 0 else WHITE)
box(s, Inches(0.30), Inches(7.03), Inches(12.73), Inches(0.08), NAVY)
notes(s, """Exercise 1.2b asks: what differences do you observe in the limited-angle case compared to Exercise 1.1?
The two degradation modes are qualitatively different:
In Ex 1.1 (reducing projections), the images get noisier but the anatomy remains geometrically correct — structures still look like the right shape, just with added noise.
In Ex 1.2 (restricting angular range), the images develop directional artefacts — structures are elongated or smeared in specific orientations corresponding to the missing wedge. At 40°, the anatomy is geometrically distorted.
This distinction matters clinically: noisy images can sometimes be denoised post-hoc; geometrically distorted images cannot — the missing information was never acquired.
GD helps in both cases but cannot recover what was never measured — the missing-wedge artefacts are fundamental, not just a noise issue.""")

# ── 15. WHY LIMITED ANGLE FAILS ──────────────────────────────────────────────
s = content_slide(CT, "Why It Fails: Missing Wedge in Fourier Space",
                  "Fourier Slice Theorem — each projection angle = one radial line in 2D Fourier domain")
box(s, Inches(0.30), Inches(1.15), Inches(6.1), Inches(5.85), LGRAY)
box(s, Inches(0.30), Inches(1.15), Inches(0.09), Inches(5.85), CT)
for i, line in enumerate([
    "Each projection at angle θ samples one radial line in Fourier space",
    "",
    "Full 0°–180° scan:  Fourier plane densely covered → accurate inversion",
    "",
    "Limited scan [0°, α]:  a wedge of Fourier space is MISSING",
    "",
    "Edges perpendicular to unsampled directions lack high-frequency info",
    "",
    "Result: blurred / elongated structures in those directions",
    "",
    "GD also fails — it cannot invent missing frequency information",
]):
    txt(s, line, Inches(0.50), Inches(1.30 + i * 0.49), Inches(5.7), Inches(0.47), size=13, color=GRAY)
pic(s, F12 + "/metrics_vs_angle.png", Inches(6.60), Inches(1.15), Inches(6.53), Inches(5.85))
notes(s, """The Fourier Slice Theorem is the theoretical foundation. It states that the 1D Fourier transform of a projection at angle θ equals a slice through the 2D Fourier transform of the image at the same angle — i.e., one radial line.
With full 180° coverage, projections from all angles are available, so the Fourier plane is densely sampled (except for some aliasing at high radii). Inversion is well-conditioned.
With only [0°, α] coverage, a wedge-shaped region of Fourier space is never sampled — frequencies in that wedge direction are simply missing from the data. You can't recover information that was never measured. GD iterates over the same underdetermined system and also cannot recover missing wedge information.
The elongation artefacts happen because the missing wedge corresponds to certain edge orientations — edges perpendicular to unsampled directions are lost. The image gets 'smeared' in those directions.
The right-hand plot shows PSNR and SSIM vs angular range for both methods — confirming that quality drops monotonically as the range decreases, and GD offers modest improvement over FBP but cannot compensate for the fundamentally missing data.""")

# ── 15b. EX 1.3 PART (a): FBP FILTERS IN THE LITERATURE ─────────────────────
s = content_slide(CT, "FBP Filter Types: Literature Review  (Exercise 1.3a)",
                  "All are windowed versions of the ramp filter — trading noise suppression for spatial resolution")
items_1_3a = [
    ("Ramp\n(Ram-Lak)",    "H(f) = |f|",
     "Full bandwidth\nSharpest edges\nMaximum noise\namplification\n\nExact inversion\nfor noiseless data"),
    ("Shepp–Logan",        "H(f) = |f|·sinc(f/2f_c)",
     "Sinc window\nMild noise reduction\nGood edge\npreservation\n\nMost common in\nclinical CT"),
    ("Cosine",             "H(f) = |f|·cos(πf/2f_c)",
     "Cosine window\nModerate rolloff\nMore smoothing\nthan Shepp-Logan"),
    ("Hann",               "H(f) = |f|·½(1+cos(πf/f_c))",
     "Raised-cosine\nwindow\nSmooth rolloff\nGood noise–\nresolution balance"),
    ("Hamming",            "H(f) = |f|·(0.54+0.46·cos(πf/f_c))",
     "Non-zero floor\nAggressive\nnoise suppression\nBest for very\nlow-dose CT"),
]
for i, (name, formula, desc) in enumerate(items_1_3a):
    x = Inches(0.28 + i * 2.57)
    w = Inches(2.45)
    box(s, x, Inches(1.15), w, Inches(5.85), LGRAY)
    box(s, x, Inches(1.15), w, Inches(0.10), CT)
    txt(s, name,    x+Inches(0.08), Inches(1.28), w-Inches(0.16), Inches(0.55), size=13, color=CT, bold=True)
    box(s, x+Inches(0.08), Inches(1.88), w-Inches(0.16), Inches(0.50), CT_L)
    txt(s, formula, x+Inches(0.10), Inches(1.91), w-Inches(0.20), Inches(0.44),
        size=9, color=CT, italic=True, align=PP_ALIGN.CENTER)
    txt(s, desc,    x+Inches(0.08), Inches(2.48), w-Inches(0.16), Inches(4.40), size=12, color=GRAY)
box(s, Inches(0.30), Inches(7.05), Inches(12.73), Inches(0.38), CT_L)
txt(s, "All windows taper the ramp at high frequencies to suppress noise — narrower passband = more smoothing = lower resolution",
    Inches(0.40), Inches(7.07), Inches(12.53), Inches(0.34),
    size=11, color=CT, italic=True, align=PP_ALIGN.CENTER)
notes(s, """Exercise 1.3 Part (a) asks for FBP filter types from the literature.
The ideal ramp filter H(f) = |f| is required for mathematically exact inversion (from the Fourier Slice Theorem), but it amplifies high-frequency noise catastrophically at low dose.
All windowed filters multiply the ramp by a smooth function that rolls off to zero near the Nyquist frequency.
Shepp-Logan: most common clinical filter. Multiplies ramp by a sinc function — mild noise reduction, good edge preservation.
Cosine and Hann: progressively more aggressive rolloff than Shepp-Logan.
Hamming: raised-cosine window with a non-zero floor (0.54 at DC), giving the smoothest rolloff and the best noise suppression — which is what I test in Part (b).
The clinical choice depends on dose and diagnostic task: high-dose scans can use the ramp; low-dose scans benefit from Shepp-Logan or Hamming.""")

# ── 15c. TABLE 3 — FBP FILTER METRICS  (Exercise 1.3b) ──────────────────────
s = content_slide(CT, "FBP Filter Comparison: Metrics  (Table 3,  Exercise 1.3b)",
                  "Low-dose setting: I₀ = 10²,  360 views  ·  Same phantom and circular mask as Ex 1.1")
t3_hdrs = ["Filter", "MSE", "PSNR (dB)", "SSIM"]
t3_rows = [
    ["Ramp  (Ram-Lak)", "0.1911", " 7.19", "0.1902"],
    ["Shepp–Logan",     "0.1250", " 9.03", "0.1951"],
    ["Hamming",         "0.0316", "15.00", "0.2291"],
]
make_table(s, t3_hdrs, t3_rows,
           [0.40, 0.20, 0.20, 0.20],
           Inches(2.50), Inches(1.80), Inches(8.33),
           gd_rows={2}, accent=CT)
for i, (filt, interp) in enumerate([
    ("Ramp",        "Sharpest but most noise — ramp amplifies every high-frequency component"),
    ("Shepp–Logan", "Moderate improvement — sinc window suppresses near-Nyquist noise slightly"),
    ("Hamming",     "Best metrics — raised-cosine window aggressively suppresses noise\n"
                    "But: strongest smoothing, fine detail is blurred → clinical trade-off"),
]):
    y = Inches(4.20 + i * 0.87)
    box(s, Inches(0.30), y, Inches(0.09), Inches(0.78), CT)
    box(s, Inches(0.39), y, Inches(12.64), Inches(0.78), LGRAY)
    txt(s, filt,   Inches(0.50), y+Inches(0.06), Inches(2.5),  Inches(0.65), size=14, color=CT, bold=True)
    txt(s, interp, Inches(3.10), y+Inches(0.06), Inches(9.80), Inches(0.65), size=13, color=GRAY)
notes(s, """Table 3 from the report: FBP filter metrics evaluated on the low-dose case from Exercise 1.1 (I₀=10², 360 views) with a fresh noise realisation.
Ramp: 7.19 dB PSNR — the mathematically 'exact' filter performs the worst in practice because it amplifies noise at all frequencies without discrimination.
Shepp-Logan: 9.03 dB — a modest improvement. The sinc window damps frequencies near Nyquist, providing slight noise reduction.
Hamming: 15.00 dB — a dramatic improvement of ~8 dB over ramp. The raised-cosine window strongly suppresses high-frequency noise.
The trade-off: Hamming gives the best scalar metrics but also the most spatial blurring. In clinical practice, if a tiny lesion needs to be resolved, the blurring introduced by Hamming could obscure it. The optimal filter depends on the diagnostic task, not just on PSNR/SSIM.""")

# ── 16. FBP FILTERS ──────────────────────────────────────────────────────────
s = content_slide(CT, "FBP Filter Comparison  (I₀ = 10²,  360 views)",
                  "Windowed ramp filters trade noise suppression for spatial resolution")
pic(s, F13 + "/exercise_1_3b_filter_comparison.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.35))
for i, (filt, psnr, note) in enumerate([
    ("Ramp (Ram-Lak)", "PSNR  7.2 dB  /  SSIM 0.190", "Sharpest — most noise"),
    ("Shepp–Logan",    "PSNR  9.0 dB  /  SSIM 0.195", "Slight noise reduction"),
    ("Hamming",        "PSNR 15.0 dB  /  SSIM 0.229", "Most smoothing — best metrics"),
]):
    x = Inches(0.35 + i * 4.34)
    box(s, x, Inches(6.58), Inches(4.14), Inches(0.78), CT_L)
    txt(s, f"{filt}   |   {psnr}   |   {note}",
        x + Inches(0.08), Inches(6.61), Inches(3.98), Inches(0.72),
        size=11, color=CT, align=PP_ALIGN.CENTER)
notes(s, """Exercise 1.3b: I reproduced the low-dose condition from 1.1 (I₀=10², 360 views) with a fresh noise realisation and compared three FBP filters.
All three are windowed versions of the ideal ramp filter, which amplifies frequency f proportionally. Windowing tapers the high frequencies to reduce noise amplification at the cost of resolution.
Ramp (Ram-Lak): no windowing — full bandwidth, sharpest edges, but the noise at low I₀ is catastrophic. PSNR 7.2 dB.
Shepp-Logan: a sinc-windowed ramp. Slight attenuation near Nyquist. Small improvement: 9.0 dB.
Hamming: a raised-cosine window that aggressively rolls off. Strong noise suppression — PSNR 15.0 dB and visibly smoother image. But fine anatomical detail is blurred.
The figure shows reconstructions and error maps side by side so you can see the trade-off visually.""")

# ── 17. NOISE vs RESOLUTION ───────────────────────────────────────────────────
s = content_slide(CT, "Noise ↔ Resolution Trade-off in Filter Choice")
box(s, Inches(0.30), Inches(1.15), Inches(12.73), Inches(5.40), LGRAY)
pic(s, F13 + "/exercise_1_3b_filter_metrics.png",
    Inches(0.35), Inches(1.20), Inches(12.63), Inches(5.30))
txt(s, "Key insight: there is no single best filter — the optimal choice depends on dose and diagnostic task",
    Inches(0.35), Inches(6.65), Inches(12.63), Inches(0.42),
    size=13, color=CT, italic=True, align=PP_ALIGN.CENTER)
notes(s, """This figure shows the MSE, PSNR, and SSIM metrics for each filter on the low-dose example.
The key tension: the ramp filter is theoretically 'exact' — it implements the precise inversion formula. But 'exact' for noiseless data becomes terrible for noisy data, because the ramp amplifies noise at high frequencies.
Smoothing filters sacrifice resolution (blur fine structures) but gain noise robustness. Hamming with PSNR 15 dB seems like a big win over Ramp's 7.2 dB. But if a radiologist needs to see a small lesion at the edge of the lung, the blurring introduced by Hamming might hide it.
In clinical practice, the filter is chosen based on the body part and dose: Shepp-Logan for standard dose chest CT, more aggressive smoothing for ultra-low dose. This is a human decision, not just a mathematical one.""")

# ── 17b. TABLES 4 + 5 — GD vs SUBSET GD METRICS  (Exercise 1.3c) ─────────────
s = content_slide(CT, "Full-Batch GD vs Subset GD: Metrics  (Tables 4 & 5,  Exercise 1.3c)",
                  "100 epochs  ·  I₀ = 10²,  360 views  ·  S = 10 subsets")
# Table 4 header
txt(s, "Table 4 — Final metrics at epoch 100  (scaled comparison,  γ_s = γ/S = 0.0001)",
    Inches(0.30), Inches(1.18), Inches(12.73), Inches(0.38), size=13, color=CT, bold=True)
t4_hdrs = ["Method", "MSE", "PSNR (dB)", "SSIM"]
t4_rows = [
    ["Full-batch GD  (γ = 0.001)",        "0.006193", "22.08", "0.3299"],
    ["Subset GD  (γ/S = 0.0001,  S = 10)", "0.006190", "22.08", "0.3300"],
]
make_table(s, t4_hdrs, t4_rows,
           [0.52, 0.16, 0.16, 0.16],
           Inches(0.30), Inches(1.60), Inches(12.73),
           gd_rows=set(), accent=CT)
# Table 5 header
txt(s, "Table 5 — Convergence summary  (epoch to peak PSNR, peak PSNR, final PSNR at epoch 100)",
    Inches(0.30), Inches(3.30), Inches(12.73), Inches(0.38), size=13, color=CT, bold=True)
t5_hdrs = ["Method", "Epoch to peak", "Peak PSNR (dB)", "Final PSNR (dB)"]
t5_rows = [
    ["Full-batch GD  (γ = 0.001)",              "33", "26.38", "22.08"],
    ["Subset GD  (γ/S = 0.0001,  S = 10)",      "34", "26.35", "22.08"],
    ["Subset GD  (γ = 0.001 per step,  S = 10)", " 3", "26.34", " 7.91"],
]
make_table(s, t5_hdrs, t5_rows,
           [0.52, 0.16, 0.16, 0.16],
           Inches(0.30), Inches(3.72), Inches(12.73),
           gd_rows={2}, accent=SEG)
box(s, Inches(0.30), Inches(5.75), Inches(6.10), Inches(0.72), LGRAY)
txt(s, "Scaled variant (γ/S): identical final quality to full-batch\n→ just a reparametrisation, same noise floor",
    Inches(0.40), Inches(5.78), Inches(5.90), Inches(0.67), size=12, color=GRAY)
box(s, Inches(6.93), Inches(5.75), Inches(6.10), Inches(0.72), RGBColor(0xff, 0xee, 0xdd))
txt(s, "Unscaled variant (γ per step): peaks at epoch 3 but diverges to 7.91 dB\n→ needs early stopping or step-size decay",
    Inches(7.03), Inches(5.78), Inches(5.90), Inches(0.67), size=12, color=SEG, bold=True)
notes(s, """Tables 4 and 5 from the report summarise the convergence experiment (Exercise 1.3c).
Table 4 (final metrics at epoch 100, scaled comparison): both methods converge to exactly the same solution — MSE 0.00619, PSNR 22.08 dB, SSIM 0.33. Scaled subset GD is mathematically equivalent to full-batch GD; the final quality is identical.
Table 5 (convergence speed): the unscaled subset GD reaches peak PSNR 26.34 dB already at epoch 3, versus epoch 33 for full-batch — 10× faster to peak. But then it diverges catastrophically: PSNR drops to 7.91 dB at epoch 100. This is because the effective per-epoch step is S×γ = 10× larger, which overshoots near the solution.
Key message: subset GD accelerates early convergence but requires careful step-size control (divide by S, or use early stopping at the peak) to be practically useful.""")

# ── 18. GD vs SUBSET GD ──────────────────────────────────────────────────────
s = content_slide(CT, "Full-Batch GD vs Subset GD  (OS-SART)",
                  "Speed–stability trade-off  ·  S = 10 subsets,  100 epochs")
pic(s, F13 + "/exercise_1_3c_psnr_trajectory.png",
    Inches(0.25), Inches(1.15), Inches(8.35), Inches(5.85))
for i, (color, title, desc) in enumerate([
    (CT,  "Full-Batch GD  (γ = 0.001)",
     "Stable convergence\nPeak PSNR 26.4 dB at epoch 33\nFinal PSNR 22.1 dB at epoch 100"),
    (MRI, "Subset GD  (γ/S = 0.0001)",
     "Equivalent final quality\nSame peak, same noise floor\nJust a reparametrisation"),
    (SEG, "Subset GD  (γ = 0.001/step)",
     "3× faster to peak (epoch 3)\nBut diverges: 7.9 dB by epoch 100\nNeeds early stopping or decay"),
]):
    y = Inches(1.20 + i * 1.95)
    box(s, Inches(8.75), y, Inches(4.35), Inches(1.80), LGRAY)
    box(s, Inches(8.75), y, Inches(0.09), Inches(1.80), color)
    txt(s, title, Inches(8.93), y + Inches(0.10), Inches(4.10), Inches(0.48), size=14, color=color, bold=True)
    txt(s, desc,  Inches(8.93), y + Inches(0.58), Inches(4.10), Inches(1.10), size=12, color=GRAY)
notes(s, """Exercise 1.3c compares full-batch GD with subset GD, which is analogous to OS-SART (ordered-subset simultaneous algebraic reconstruction technique).
What I did: split the 360 projections into S=10 interleaved subsets of 36 projections each. Within each epoch, I update the image once per subset sequentially — so each epoch makes 10 update steps, each using 1/10 of the data.
Two step-size choices:
  1. γ_s = γ/S = 0.0001 (scaled): each sub-step is 10× smaller, so the total update per epoch matches full-batch GD. Result: identical convergence curves — same peak PSNR 26.4 dB at epoch 33-34, same final PSNR 22.1 dB. This is mathematically predictable.
  2. γ_s = γ = 0.001 (unscaled): each sub-step is the same size as a full-batch step, giving 10× more total update per epoch. Result: it reaches peak PSNR 26.3 dB at epoch 3 — 10× faster! But then it diverges catastrophically; by epoch 100 PSNR is only 7.9 dB. This is because the effective learning rate is too large once we're near the solution.
The lesson: subset GD accelerates early convergence but requires early stopping or step-size decay to be practical.""")

# ── 19. PET vs CT ────────────────────────────────────────────────────────────
s = content_slide(CT, "Why PET Uses MLEM,  CT Uses Gradient Descent")
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(5.85), RGBColor(0xff, 0xf8, 0xe1))
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(0.10), SEG)
txt(s, "PET", Inches(0.48), Inches(1.30), Inches(5.65), Inches(0.55), size=22, color=SEG, bold=True)
for i, p in enumerate(["Measurements = photon counts",
                        "→ naturally Poisson-distributed",
                        "MLEM: Maximum-Likelihood EM algorithm",
                        "Directly models Poisson likelihood",
                        "ℓ(x) = Σ [yᵢ log(Ax)ᵢ − (Ax)ᵢ]",
                        "Classic: Shepp & Vardi formulation",
                        "Also: OSEM (ordered-subset EM)"]):
    txt(s, "•  " + p, Inches(0.48), Inches(1.97 + i * 0.68), Inches(5.65), Inches(0.62), size=13, color=GRAY)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(5.85), CT_L)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(0.10), CT)
txt(s, "CT", Inches(7.25), Inches(1.30), Inches(5.75), Inches(0.55), size=22, color=CT, bold=True)
for i, p in enumerate(["Measurements = transmission ratios",
                        "→ log-transform → approx. Gaussian",
                        "Gradient Descent / SART / SIRT",
                        "Minimises least-squares objective",
                        "‖Ax − b‖²   (Gaussian data fit)",
                        "SART: algebraic row-action update",
                        "GPU-accelerated in clinical systems"]):
    txt(s, "•  " + p, Inches(7.25), Inches(1.97 + i * 0.68), Inches(5.75), Inches(0.62), size=13, color=GRAY)
box(s, Inches(6.0), Inches(3.55), Inches(1.13), Inches(0.90), WHITE)
txt(s, "Data\nmodel\n↓", Inches(6.0), Inches(3.55), Inches(1.13), Inches(0.90),
    size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, """Exercise 1.3d asks: why does PET use MLEM while CT uses gradient-descent-type methods?
The answer is: the choice of optimisation algorithm should match the statistical model of the noise.
In PET: you inject a radiotracer and detect pairs of photons from positron annihilation. What you directly count is photons — and photon counts are inherently Poisson-distributed. MLEM (Maximum-Likelihood Expectation-Maximisation) is the algorithm that maximises the Poisson likelihood. It was derived specifically for emission tomography by Shepp and Vardi. Each MLEM update is multiplicative and always keeps the image non-negative, which matters since activity cannot be negative.
In CT: the detector measures transmitted X-ray intensity I. You then apply a log transform: b = -log(I/I₀) · p_max. After this transformation, the noise is approximately Gaussian (for reasonably high I₀). Gaussian noise → least-squares objective → gradient descent. SART and SIRT are algebraic variants of the same least-squares idea.
The key principle: always match your reconstruction algorithm to the noise model of your measurements.""")

# ── 20. MRI SECTION HEADER ───────────────────────────────────────────────────
s = section_slide(MRI, MRI_BADGE, "MODULE  2",
                  "MRI Denoising",
                  "k-space  ·  Multi-coil RSS  ·  Image-space denoising  ·  k-space filtering")
notes(s, """Moving to Module 2: MRI.
This module is purely practical — loading real knee MRI data, understanding what k-space looks like, reconstructing images from complex-valued multi-coil data, and then comparing several denoising approaches.
No forward model simulation here — we work with measured data.""")

# ── 21. K-SPACE ──────────────────────────────────────────────────────────────
s = content_slide(MRI, "MRI: Data Acquired in k-Space",
                  "6 receiver coils  ·  log|k(u,v)| shown  ·  Knee dataset (280×280, complex128)")
pic(s, F21 + "/kspace_magnitude_all_coils.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.55))
box(s, Inches(0.20), Inches(6.78), Inches(6.1), Inches(0.56), MRI_L)
txt(s, "Bright centre = low frequencies = global image structure",
    Inches(0.30), Inches(6.80), Inches(5.9), Inches(0.52), size=13, color=MRI, bold=True)
box(s, Inches(6.70), Inches(6.78), Inches(6.43), Inches(0.56), MRI_L)
txt(s, "Coil-to-coil differences reflect spatially varying receive sensitivity",
    Inches(6.80), Inches(6.80), Inches(6.23), Inches(0.52), size=13, color=MRI, bold=True)
notes(s, """What I did: loaded the file knee.npy which contains a 6×280×280 complex128 array — 6 coils, each with a 280×280 complex k-space matrix.
In MRI, the scanner doesn't acquire an image directly. It acquires the 2D Fourier transform of the tissue magnetisation — called k-space. Each k-space point k(u,v) corresponds to one spatial frequency component of the image.
I visualised log(1+|k|) for each coil — the log scale is necessary because the dynamic range is enormous: the DC component (centre, u=v=0) is thousands of times larger than the outer high-frequency components.
The bright central cross/peak corresponds to low spatial frequencies — these encode the bulk of the image contrast, the large smooth regions. The dim outer ring contains high-frequency content — fine edges and small features. Noise lives here too, which is why filtering in k-space is equivalent to low-pass filtering the image.
Coil-to-coil differences: each physical coil element has a different receive sensitivity profile, so the amplitude pattern differs between coils.""")

# ── 22. RECONSTRUCTION ───────────────────────────────────────────────────────
s = content_slide(MRI, "k-Space → Image Space  (2D IFFT)",
                  "Complex image: magnitude reveals anatomy  ·  phase shows slow field variation")
pic(s, F21 + "/magnitude_phase_coil0.png",
    Inches(0.40), Inches(1.15), Inches(12.53), Inches(5.55))
box(s, Inches(0.35), Inches(6.78), Inches(5.5), Inches(0.56), MRI_L)
txt(s, "Magnitude  —  knee cartilage and cortical bone clearly visible",
    Inches(0.45), Inches(6.80), Inches(5.3), Inches(0.52), size=13, color=MRI, bold=True)
box(s, Inches(7.30), Inches(6.78), Inches(5.8), Inches(0.56), CT_L)
txt(s, "Phase  —  slowly varying field with wrap transitions  (less diagnostic)",
    Inches(7.40), Inches(6.80), Inches(5.6), Inches(0.52), size=13, color=CT, bold=True)
notes(s, """To reconstruct an image from k-space, I applied np.fft.ifft2 to each coil's k-space — a 2D inverse DFT.
The result is a complex-valued image I(x,y) = M(x,y) · exp(iφ(x,y)).
Magnitude |I(x,y)|: this is the anatomical image. For coil 0 you can see the knee clearly — the femoral condyle, tibial plateau, cartilage layers, and cortical bone margins.
Phase angle(I(x,y)): this reflects the B0 field inhomogeneity and receiver coil phase effects. It shows a slowly varying gradient with wrap-around transitions at ±π. This is not useful for anatomical interpretation — it's a confound from the acquisition, not a tissue property.
In clinical MRI, we almost always use magnitude images. Phase images are used for specific applications like flow imaging or susceptibility-weighted imaging.""")

# ── 23. MULTI-COIL ───────────────────────────────────────────────────────────
s = content_slide(MRI, "6 Receiver Coils: Different Spatial Sensitivity",
                  "Same anatomy resolved — brightness distribution differs across coils")
pic(s, F21 + "/magnitude_all_coils.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.58))
box(s, Inches(0.25), Inches(6.82), Inches(12.83), Inches(0.52), MRI_L)
txt(s, "Each coil is most sensitive to nearby tissue  →  complementary coverage  →  "
       "no single coil gives uniform illumination of the full FOV",
    Inches(0.35), Inches(6.84), Inches(12.63), Inches(0.48),
    size=13, color=MRI, bold=True, align=PP_ALIGN.CENTER)
notes(s, """This shows all six coil magnitude images together.
All six resolve the same knee anatomy — the bone, cartilage, and soft tissue structures are consistent.
But the brightness distribution is completely different. Coil 0: bright in the upper-right. Coil 3: bright in the lower-left. Coil 5: bright at the bottom with a dark region at top. This is because each physical receive coil element is positioned at a different location around the scanner bore, and its SNR drops off with distance from the element.
The result is that no individual coil gives uniform illumination. Regions far from a coil appear dim (low SNR) in that coil's image. This is why we need coil combination — each coil contributes good SNR in a different part of the field of view.""")

# ── 24. RSS ──────────────────────────────────────────────────────────────────
s = content_slide(MRI, "Root-Sum-of-Squares (RSS) Combination",
                  "I_RSS(x,y) = √Σ|I_c(x,y)|²   ·   Combines 6 coils into one uniformly illuminated image")
pic(s, F21 + "/rss_combined.png",
    Inches(2.65), Inches(1.10), Inches(8.03), Inches(6.12))
txt(s, "Regions dim\nin every\nindividual coil\nare now\nclearly visible",
    Inches(0.15), Inches(3.10), Inches(2.40), Inches(1.80),
    size=14, color=MRI, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Full knee\ncross-section\nuniformly\nilluminated\nacross FOV",
    Inches(10.78), Inches(3.10), Inches(2.40), Inches(1.80),
    size=14, color=MRI, bold=True, align=PP_ALIGN.CENTER)
notes(s, """The RSS combination formula: at each pixel (x,y), take the square root of the sum of squared magnitudes across all 6 coils.
This is the standard combination method used in virtually all modern clinical MRI scanners. It's statistically justified: under the assumption of independent Gaussian noise across coils, RSS is approximately the maximum-likelihood estimate of the signal amplitude.
The result shows clear improvement: regions that were dark in every individual coil (because they were far from all coil elements) are now visible. The full knee cross-section is uniformly bright. SNR is improved across the entire field of view compared to any single coil.
One limitation: RSS doesn't account for noise correlation between coils. In practice, matched filter combination (also called SENSE combination) is optimal, but it requires knowing each coil's sensitivity map, which I didn't estimate here.""")

# ── 24b. EXERCISE 2.1 OBSERVATIONS ──────────────────────────────────────────
s = content_slide(MRI, "Exercise 2.1: Observations",
                  "Key findings from k-space visualisation, image reconstruction, and coil combination")
obs = [
    ("k-Space structure",
     "All six coils show the characteristic bright central peak (low frequencies = global image contrast) "
     "surrounded by dim outer regions (high frequencies = edges and fine detail). "
     "Noise is distributed across all frequencies but dominates at the periphery."),
    ("Coil-to-coil differences",
     "The magnitude of the central peak and the fall-off pattern differ between coils, "
     "reflecting each coil element's spatially varying receive sensitivity. "
     "No single coil uniformly covers the full FOV."),
    ("Magnitude vs phase",
     "The magnitude image reveals clear anatomical structure (knee cartilage, cortical bone, soft tissue). "
     "The phase image shows a slowly varying background field with ±π wrap transitions "
     "— a confound from the acquisition, not a tissue property."),
    ("RSS combination benefit",
     "The RSS image is uniformly illuminated across the full knee cross-section. "
     "Regions that were dim in every individual coil are now clearly visible. "
     "This confirms that multi-coil acquisition provides complementary coverage."),
]
for i, (title, desc) in enumerate(obs):
    y = Inches(1.25 + i * 1.42)
    box(s, Inches(0.30), y, Inches(0.09), Inches(1.28), MRI)
    box(s, Inches(0.39), y, Inches(12.64), Inches(1.28), LGRAY)
    txt(s, title, Inches(0.52), y + Inches(0.08), Inches(3.0),  Inches(0.45), size=14, color=MRI, bold=True)
    txt(s, desc,  Inches(0.52), y + Inches(0.52), Inches(12.3), Inches(0.68), size=13, color=GRAY)
notes(s, """This slide corresponds to the 'Observations' section required by Exercise 2.1.
Four key observations:
1. k-space structure: bright centre (low freq = bulk signal), dim periphery (high freq = edges+noise). Log scaling is essential to visualise both.
2. Coil sensitivity: each coil sees the same anatomy but with different spatial weighting — coil nearest the tissue gives brightest signal in that region.
3. Magnitude is diagnostically useful; phase mainly reflects field inhomogeneity and has limited anatomical value in standard MRI.
4. RSS combination solves the non-uniform coverage problem — it combines the complementary sensitivity profiles of all 6 coils into one uniformly bright image.""")

# ── 25. IMAGE-SPACE DENOISING ────────────────────────────────────────────────
s = content_slide(MRI, "Image-Space Denoising: Three Methods  (Coil 0)",
                  "Applied independently to each coil magnitude image before RSS combination")
pic(s, F22 + "/part_a_comparison_coil0.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.05))
methods = [
    ("Original",      "Local std: 0.038",  "Noisy baseline",               LGRAY, GRAY),
    ("Median 3×3",    "Local std: 0.016",  "Reduces speckle, keeps edges",  LGRAY, GRAY),
    ("Gaussian σ=1",  "Local std: 0.011",  "Blurs anatomical boundaries",   LGRAY, GRAY),
    ("Bilateral",     "Local std: 0.004",  "Best: noise ↓, edges preserved", MRI,  WHITE),
]
for i, (name, metric, note, bg_c, tc) in enumerate(methods):
    x = Inches(0.20 + i * 3.30)
    box(s, x, Inches(6.28), Inches(3.20), Inches(1.08), bg_c)
    txt(s, f"{name}\n{metric}  ·  {note}",
        x+Inches(0.08), Inches(6.31), Inches(3.04), Inches(1.02),
        size=11, color=tc, bold=(tc==WHITE), align=PP_ALIGN.CENTER)
notes(s, """Exercise 2.2a: I applied three denoising filters directly to the magnitude images of each coil BEFORE the RSS combination step.
The three methods:
  — Median filter (3×3 window): replaces each pixel with the median of its 3×3 neighbourhood. Reduces salt-and-pepper noise. Local std drops from 0.038 to 0.016. Keeps edges slightly sharper than Gaussian because it doesn't average across edge boundaries.
  — Gaussian filter (σ=1.0 px): convolves with a Gaussian kernel. Reduces noise uniformly. Local std 0.011. But it blurs the cartilage boundary and cortical bone margins visibly — anywhere there's a sharp transition, Gaussian smooths it.
  — Bilateral filter (σ_spatial=3.0 px, σ_colour=0.05): a joint spatial/intensity filter. It weights neighbours by both spatial distance AND intensity similarity. So pixels across a sharp edge (different intensity) get low weight — the edge is preserved. Flat regions denoise strongly. Local std 0.004 — the best noise reduction. Visually, the cartilage boundary remains crisp.
I chose bilateral as the recommended method for the final pipeline.""")

# ── 25b. ALL COILS — MEDIAN FILTERING ────────────────────────────────────────
s = content_slide(MRI, "Median Filtering (3×3)  —  All Six Coils  (Fig. 20)",
                  "Consistent noise reduction across all coils  ·  Thin structures slightly flattened")
pic(s, F22 + "/part_a_all_coils_median.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
notes(s, """Figure 20 from the report: all six coil magnitude images after median filtering with a 3×3 window.
The median filter is consistent across coils: speckle noise is reduced in all six, including the lower-SNR coils far from the anatomy.
Compared to Gaussian: edges are preserved slightly better because the median operation does not average across discontinuities — it picks the median value in the neighbourhood, which is more likely to be on one side of the edge.
Compared to bilateral: thinner structures can appear flattened when the 3×3 window spans across them. Local std drops from 0.038 (original) to 0.016.""")

# ── 25c. ALL COILS — GAUSSIAN FILTERING ──────────────────────────────────────
s = content_slide(MRI, "Gaussian Filtering (σ=1.0 px)  —  All Six Coils  (Fig. 21)",
                  "Uniform smoothing across coils  ·  Strongest blurring of anatomical boundaries")
pic(s, F22 + "/part_a_all_coils_gaussian.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
notes(s, """Figure 21 from the report: all six coil magnitude images after Gaussian filtering with σ=1.0 px.
Gaussian filtering produces the most uniform noise reduction — every region is convolved with the same Gaussian kernel regardless of local content.
The downside is clear in the higher-SNR coils: the cartilage boundary is visibly softer, and fine bone detail is blurred. The filter cannot distinguish signal edges from noise — it treats all high-frequency content identically.
Local std drops to 0.011. The strong noise reduction is achieved at the cost of spatial resolution, which can hide clinically relevant fine structures.""")

# ── 25d. ALL COILS — BILATERAL FILTERING ─────────────────────────────────────
s = content_slide(MRI, "Bilateral Filtering (σ_s=3.0, σ_c=0.05)  —  All Six Coils  (Fig. 22)",
                  "Best noise reduction  ·  Main anatomical boundaries preserved  ·  Preferred method")
pic(s, F22 + "/part_a_all_coils_bilateral.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
box(s, Inches(0.25), Inches(7.05), Inches(12.83), Inches(0.38), MRI_L)
txt(s, "Local std = 0.004  —  best of all three  ·  Bilateral chosen for the final denoise-then-combine pipeline",
    Inches(0.35), Inches(7.07), Inches(12.63), Inches(0.34),
    size=12, color=MRI, bold=True, align=PP_ALIGN.CENTER)
notes(s, """Figure 22 from the report: all six coil images after bilateral filtering (σ_spatial=3.0 px, σ_colour=0.05).
Bilateral achieves the best noise reduction (local std 0.004) while preserving the main anatomical boundaries.
The intensity-weighting means pixels across a sharp edge (with very different intensity) get low weight in the smoothing — the edge is not blurred. Flat regions denoise strongly while edges are protected.
In the noisiest coils, bilateral can sometimes suppress weak texture that is actually signal, but for this dataset the balance is good.
This is the filter chosen for the final denoise-then-combine pipeline in Part (c).""")

# ── 25e. TABLE 6 — DENOISING LOCAL STD  (Exercise 2.2a) ─────────────────────
s = content_slide(MRI, "Image-Space Denoising: Quantitative Summary  (Table 6,  Exercise 2.2a)",
                  "Local std measured in a fixed 20×20 low-signal region of coil 0  (rows 250–270, cols 30–50)")
t6_hdrs = ["Method", "Local std  (noise proxy)"]
t6_rows = [
    ["Original magnitude",  "0.03822"],
    ["Median  (3×3)",       "0.01579"],
    ["Gaussian  (σ = 1.0)", "0.01059"],
    ["Bilateral",           "0.00390"],
]
make_table(s, t6_hdrs, t6_rows,
           [0.60, 0.40],
           Inches(2.80), Inches(1.80), Inches(7.73),
           gd_rows={3}, accent=MRI)
txt(s, "Note: local std is a simple noise proxy — not a full image-quality metric (no ground truth available)",
    Inches(0.30), Inches(4.22), Inches(12.73), Inches(0.38),
    size=12, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
for i, (method, reduction, note) in enumerate([
    ("Median 3×3",    "58.7% reduction",  "Moderate — good edge preservation, removes speckle"),
    ("Gaussian σ=1",  "72.3% reduction",  "Strong — but blurs anatomical boundaries uniformly"),
    ("Bilateral",     "89.8% reduction",  "Best — edge-preserving; smooths background strongly"),
]):
    y = Inches(4.75 + i * 0.78)
    box(s, Inches(0.30), y, Inches(0.09), Inches(0.70), MRI)
    box(s, Inches(0.39), y, Inches(12.64), Inches(0.70), LGRAY)
    txt(s, method,    Inches(0.50), y+Inches(0.06), Inches(2.5), Inches(0.58), size=13, color=MRI, bold=True)
    txt(s, reduction, Inches(3.10), y+Inches(0.06), Inches(2.3), Inches(0.58), size=13, color=NAVY, bold=True)
    txt(s, note,      Inches(5.50), y+Inches(0.06), Inches(7.4), Inches(0.58), size=13, color=GRAY)
notes(s, """Table 6 from the report: standard deviation in a fixed 20×20 low-signal region of coil 0.
This is a simple noise proxy — we measure residual fluctuations in a region that should be near-zero signal. Lower std = more noise suppressed.
Original: 0.038 (noisy baseline).
Median 3×3: 0.016 — 59% noise reduction with moderate edge preservation.
Gaussian σ=1.0: 0.011 — 72% reduction but strongest blurring of boundaries.
Bilateral: 0.004 — 90% reduction while preserving main edges. Best of all three.
Important caveat: no ground-truth clean image is available for MRI (unlike CT where we had the original phantom). So PSNR/SSIM cannot be computed — the local std is the best simple quantitative comparison available.""")

# ── 26. K-SPACE FILTERING ────────────────────────────────────────────────────
s = content_slide(MRI, "k-Space Low-Pass Filtering  (Butterworth)",
                  "Applied in k-space before IFFT  ·  Affects both magnitude AND phase")
pic(s, F22 + "/part_b_kspace_filter.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.45))
box(s, Inches(0.25), Inches(6.68), Inches(12.83), Inches(0.68), MRI_L)
txt(s, "Butterworth LP filter (D₀=30 px, order=2) attenuates outer k-space  →  "
       "smoother image but blurs fine anatomy  ·  "
       "Global operation: no local edge adaptation",
    Inches(0.35), Inches(6.71), Inches(12.63), Inches(0.62),
    size=13, color=MRI, align=PP_ALIGN.CENTER)
notes(s, """Exercise 2.2b: instead of denoising in image space, I filtered in k-space before reconstruction.
I constructed a Butterworth low-pass filter mask in k-space coordinates. The Butterworth filter has the form H(d) = 1 / (1 + (d/D₀)^(2n)) where d is the pixel distance from the k-space centre, D₀=30 px is the cutoff, and n=2 is the order.
I multiplied the k-space of coil 0 by this mask, then applied IFFT to reconstruct the filtered image.
Key difference from image-space filtering: the k-space filter is a global operation. It multiplies every frequency component by the same weight H(d), which is equivalent to convolving the image with a low-pass kernel. There's no local adaptation — it treats edges and flat regions identically.
Result: the magnitude image is smoother, but fine anatomical detail (fine bone trabeculae, thin cartilage margins) is blurred. The phase image also changes — k-space filtering affects the complex-valued data before magnitude extraction, so both magnitude and phase are modified.""")

# ── 26b. BUTTERWORTH — MAGNITUDE AND PHASE (Fig. 24) ─────────────────────────
s = content_slide(MRI, "Butterworth: Effect on Magnitude and Phase  (Fig. 24)",
                  "Pre-reconstruction k-space filter modifies both magnitude and phase")
pic(s, F22 + "/part_b_magnitude_phase.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
notes(s, """Figure 24 from the report: coil 0 magnitude (top row) and phase (bottom row), original left vs Butterworth-filtered right.
Magnitude: the filtered image is smoother — noise suppressed but fine anatomical detail (thin bone, cartilage margins) also blurred.
Phase: the phase map is also changed by the Butterworth filter. This is the fundamental difference from image-space denoising — because we multiply the complex k-space before IFFT, both the real and imaginary parts change, which means both magnitude and phase of the reconstructed image are affected.
Image-space filters (median, Gaussian, bilateral) operate only on the derived magnitude image — they cannot change the phase.
For standard anatomical MRI this phase change is unimportant, but it matters for applications that use phase (flow imaging, susceptibility-weighted imaging, quantitative MRI).""")

# ── 27. DOMAIN COMPARISON ────────────────────────────────────────────────────
s = content_slide(MRI, "Where to Denoise: k-Space vs Image Space")
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(5.85), MRI_L)
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(0.10), MRI)
txt(s, "k-Space Filtering", Inches(0.48), Inches(1.30), Inches(5.65), Inches(0.52), size=20, color=MRI, bold=True)
for i, p in enumerate(["•  Applied before IFFT",
                        "•  Global smoothing (all frequencies equally weighted)",
                        "•  Affects both magnitude and phase",
                        "•  Natural fit for MR acquisition model",
                        "•  Closest to Gaussian in image space",
                        "•  No local edge adaptation"]):
    txt(s, p, Inches(0.48), Inches(1.97 + i * 0.68), Inches(5.65), Inches(0.62), size=13, color=GRAY)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(5.85), CT_L)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(0.10), CT)
txt(s, "Image-Space Filtering", Inches(7.25), Inches(1.30), Inches(5.75), Inches(0.52), size=20, color=CT, bold=True)
for i, p in enumerate(["•  Applied after IFFT (on magnitude)",
                        "•  Local: adapts to spatial context and edges",
                        "•  Edge-preserving (bilateral filter)",
                        "•  Operates only on derived magnitude",
                        "•  Better visual quality — preserves anatomy",
                        "•  More surgical noise suppression"]):
    txt(s, p, Inches(7.25), Inches(1.97 + i * 0.68), Inches(5.75), Inches(0.62), size=13, color=GRAY)
box(s, Inches(5.53), Inches(6.45), Inches(2.27), Inches(0.55), MRI)
txt(s, "→  Bilateral wins for MRI",
    Inches(5.53), Inches(6.48), Inches(2.27), Inches(0.50),
    size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
notes(s, """Comparing the two domains conceptually:
k-space filtering: you're operating in Fourier space, so the filter is a frequency-domain multiplier. This is efficient and elegant, but it's inherently global — every pixel in the output is influenced by every k-space sample. You can't selectively preserve an edge while smoothing its neighbours, because edge information is distributed across all frequencies.
Image-space filtering: you work on the reconstructed image directly. Local filters like median and bilateral look at a neighbourhood around each pixel. Bilateral filtering explicitly compares pixel intensities: if two nearby pixels have very different intensities (like pixels on either side of an edge), they get low weight in the smoothing — the edge is not blurred.
For MRI where preserving anatomical boundaries (cartilage margins, lesion edges) is clinically critical, image-space bilateral filtering wins. In applications where you specifically want to suppress high-frequency acquisition artefacts that live in k-space, k-space filtering can be more targeted.""")

# ── 27b. PART (c): DENOISE-THEN-COMBINE — BILATERAL RESULT (Fig. 25) ─────────
s = content_slide(MRI, "Denoise-Then-Combine: Bilateral Pipeline  (Part c,  Fig. 25)",
                  "Original RSS  ·  Bilateral-denoised then RSS  ·  Absolute pixel-wise difference")
pic(s, F22 + "/part_c_combined.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.85))
notes(s, """Figure 25 from the report: the three-panel comparison for the bilateral denoise-then-combine pipeline.
Left: original RSS combined image — the noisy baseline.
Centre: bilateral-denoised-then-combined — each of the 6 coil magnitude images was denoised with the bilateral filter independently, then all six were combined with RSS.
Right: absolute pixel-wise difference — shows where denoising made the biggest change. Bright (red) regions correspond to areas where noise was suppressed; dark regions where signal was already clean.
The centre panel is visibly smoother than the left while retaining major anatomical structure. The difference map confirms that denoising concentrated in noisy texture regions, not at the main anatomical boundaries.""")

# ── 28. FINAL PIPELINE ───────────────────────────────────────────────────────
s = content_slide(MRI, "Final Pipeline: Bilateral Denoise → RSS",
                  "Each coil denoised independently then combined  ·  Bilateral gives best balance")
pic(s, F22 + "/part_c_all_methods.png",
    Inches(0.20), Inches(1.15), Inches(12.93), Inches(5.40))
methods2 = [
    ("Original RSS",    "Noisy baseline",          LGRAY, GRAY),
    ("Median → RSS",   "Moderate denoising",       LGRAY, GRAY),
    ("Gaussian → RSS", "Strongest blur",           LGRAY, GRAY),
    ("Bilateral → RSS","Best — noise ↓, edges ✓", MRI,  WHITE),
]
for i, (name, qual, bg_c, tc) in enumerate(methods2):
    x = Inches(0.20 + i * 3.30)
    box(s, x, Inches(6.63), Inches(3.22), Inches(0.72), bg_c)
    txt(s, f"{name}  ·  {qual}",
        x+Inches(0.07), Inches(6.65), Inches(3.08), Inches(0.68),
        size=11, color=tc, bold=(tc==WHITE), align=PP_ALIGN.CENTER)
notes(s, """The final step in Exercise 2.2c: apply the best denoising method to all six coil images individually, then combine with RSS.
Pipeline: for each of the 6 coils, apply bilateral filter → take magnitude → RSS combination across all 6.
This figure shows all four outputs on the same colour scale so the comparison is fair.
Left to right: original RSS (noisy), median then RSS (moderate), Gaussian then RSS (over-smoothed), bilateral then RSS (best).
The bilateral result is visibly the cleanest: the background noise is suppressed, the cartilage boundary remains sharp, and the tibial cartilage layer is still distinguishable. The Gaussian result over-blurs and the cartilage boundary is lost. The median result is a reasonable middle ground but bilateral still wins.
I recommend bilateral as the final denoising step for this knee MRI dataset.""")

# ── 28b. DISCUSSION: COMPARISON OF ALL DENOISING METHODS (Ex 2.2) ────────────
s = content_slide(MRI, "Discussion: Comparison of All Denoising Methods  (Ex 2.2)",
                  "Image-space methods vs k-space filtering — strengths, weaknesses, and context")
comparisons = [
    (MRI,  "Bilateral filter",
     "Best visual result. Edge-preserving: spatial + intensity weighting means sharp "
     "anatomical boundaries (cartilage, cortical bone) are not blurred. "
     "Strongest noise reduction (local std 0.004). "
     "Can over-smooth weak signal in very low-SNR coils. "
     "→ Recommended for this dataset."),
    (GRAY, "Median filter (3×3)",
     "Moderate noise reduction (std 0.016). Removes speckle while keeping edges slightly "
     "sharper than Gaussian. Thin structures can be flattened when the window spans them. "
     "Good middle-ground option when bilateral is too slow."),
    (GRAY, "Gaussian filter (σ=1.0)",
     "Uniform smoothing, strongest blurring of boundaries (std 0.011). "
     "Equivalent to a low-pass filter in image space — does not distinguish "
     "signal edges from noise. Can hide clinically relevant fine detail."),
    (CT,   "Butterworth k-space filter (D₀=30, n=2)",
     "Global low-pass filter applied before IFFT — attenuates outer k-space (high frequencies). "
     "Closest to Gaussian smoothing: reduces noise but blurs fine anatomy. "
     "Unique: affects both magnitude AND phase (pre-reconstruction). "
     "No local edge adaptation. Best suited to suppressing acquisition artefacts "
     "that are localised in k-space rather than diffuse noise."),
]
for i, (color, title, desc) in enumerate(comparisons):
    y = Inches(1.18 + i * 1.45)
    box(s, Inches(0.30), y, Inches(0.09), Inches(1.30), color)
    box(s, Inches(0.39), y, Inches(12.64), Inches(1.30),
        RGBColor(0xe8, 0xf5, 0xe9) if i == 0 else LGRAY if i % 2 == 0 else WHITE)
    txt(s, title, Inches(0.52), y + Inches(0.07), Inches(3.0),  Inches(0.45), size=14, color=color, bold=True)
    txt(s, desc,  Inches(0.52), y + Inches(0.52), Inches(12.3), Inches(0.70), size=12, color=GRAY)
notes(s, """This is the 'Compare and discuss the different denoising methods' section required by Exercise 2.2.
Key comparison across all four methods:
Bilateral wins overall: it is the only method that actively preserves edges while suppressing noise, because it weights neighbours by both spatial distance AND intensity similarity. For MRI where preserving cartilage margins and lesion boundaries is clinically critical, bilateral is the best choice.
Median is a good compromise: cheaper to compute than bilateral, better edge preservation than Gaussian, but can flatten thin structures.
Gaussian gives the strongest uniform smoothing but also the most blurring — bad for fine anatomical detail.
Butterworth k-space filter is fundamentally different: it operates before reconstruction, affects phase as well as magnitude, and is a global operation with no local adaptation. Its effect in the image domain is closest to Gaussian smoothing. It's most useful for suppressing specific k-space artefacts, not diffuse noise.
The key takeaway: image-space methods (especially bilateral) are better for preserving diagnostic anatomy; k-space filtering is more appropriate for structured artefacts.""")

# ── 28c. FURTHER IMPROVEMENTS (Ex 2.2) ───────────────────────────────────────
s = content_slide(MRI, "Further Improvements to Combined Image Quality  (Ex 2.2)",
                  "Methods that could go beyond median / Gaussian / bilateral for this knee MRI dataset")
improvements = [
    ("Non-local Means (NLM)",
     "Instead of averaging only spatial neighbours, NLM searches the entire image for patches "
     "with similar structure and averages them. "
     "Exploits self-similarity in MRI (e.g. repeated bone texture). "
     "Preserves fine detail better than bilateral at the cost of higher computation."),
    ("Wavelet Denoising",
     "Transform each coil image to the wavelet domain, threshold small coefficients "
     "(noise concentrates at fine-scale wavelet levels), invert. "
     "Separates signal (large coefficients at coarse scales) from noise (small at fine scales). "
     "Naturally multi-scale — adapts to image content across spatial frequencies."),
    ("Sensitivity-Weighted Coil Combination",
     "Replace plain RSS with a weighted sum using each coil's estimated sensitivity map: "
     "I_combined = Σ s_c* · I_c / Σ |s_c|².  "
     "This is the optimal (maximum-SNR) combination and accounts for noise correlation between coils. "
     "Requires estimating sensitivity maps from a calibration scan or from the data itself."),
    ("Deep Learning Denoiser (e.g. DnCNN, MRI-specific U-Net)",
     "Train a CNN on pairs of noisy/clean knee MRI images. "
     "Learns the noise statistics of this specific scanner/protocol. "
     "Can outperform handcrafted filters significantly at inference time. "
     "Requires a training dataset, but pre-trained models for knee MRI exist."),
]
for i, (title, desc) in enumerate(improvements):
    y = Inches(1.18 + i * 1.45)
    box(s, Inches(0.30), y, Inches(0.09), Inches(1.30), MRI)
    box(s, Inches(0.39), y, Inches(12.64), Inches(1.30), LGRAY if i % 2 == 0 else WHITE)
    txt(s, title, Inches(0.52), y + Inches(0.07), Inches(12.3), Inches(0.40), size=14, color=MRI, bold=True)
    txt(s, desc,  Inches(0.52), y + Inches(0.50), Inches(12.3), Inches(0.72), size=12, color=GRAY)
notes(s, """Exercise 2.2 asks: 'Discuss any methods that might be able to further improve the combined image quality.'
Four directions:
1. Non-local means: the key advantage over bilateral is that it finds similar patches globally, not just in a small spatial neighbourhood. Very effective for MRI where tissue regions have self-similar texture. The main cost is computation (searching the full image for matches).
2. Wavelet denoising: treats noise as high-frequency content in the wavelet domain. Threshold-based: coefficients below a threshold are set to zero (hard thresholding) or shrunk (soft thresholding). The BM3D algorithm combines wavelet-like transforms with non-local grouping and is state-of-the-art for Gaussian denoising.
3. Sensitivity-weighted combination: the real limitation of RSS is that it does not account for the noise correlation between coils or their sensitivity profiles. SENSE and matched filter combination are optimal — they maximise SNR across the full FOV. This would require estimating sensitivity maps, but the improvement in SNR and uniformity can be substantial.
4. Deep learning: modern MRI denoising networks (e.g. MoDL, E2E-VarNet) are trained end-to-end on large datasets and learn the full noise model of the acquisition. They substantially outperform handcrafted filters on benchmark datasets, but require training data specific to the scanner and protocol.""")

# ── 29. SEGMENTATION SECTION HEADER ─────────────────────────────────────────
s = section_slide(SEG, SEG_BADGE, "MODULE  3",
                  "Segmentation",
                  "Deep learning (R2U-Net)  vs  Atlas-based methods")
notes(s, """Module 3 is a literature review module — no code, but critical analysis of two segmentation papers.
I selected one data-driven DL paper and one traditional atlas-based paper, then compared them and proposed improvements to the DL approach.""")

# ── 30. SEGMENTATION PROBLEM ─────────────────────────────────────────────────
s = content_slide(SEG, "Segmentation: Pixel-Wise Labelling of Structures",
                  "Assign each pixel a structure label  ·  Reduces expert annotation burden")
steps = [
    ("Input",    "Medical image\n(CT / MRI / PET)", "Raw acquisition"),
    ("Task",     "Pixel classification\n& ROI delineation", "Each pixel → organ / tumour / vessel"),
    ("Output",   "Binary or\nmulti-class mask", "Delineated structure"),
    ("Challenge","Ill-posed + noisy\n+ class imbalance", "Many images map to same mask"),
]
for i, (step, title, desc) in enumerate(steps):
    x = Inches(0.25 + i * 3.25)
    box(s, x, Inches(1.45), Inches(3.08), Inches(4.85), LGRAY)
    box(s, x, Inches(1.45), Inches(3.08), Inches(0.10), SEG)
    txt(s, step,  x+Inches(0.12), Inches(1.60), Inches(2.84), Inches(0.42), size=13, color=SEG, bold=True)
    txt(s, title, x+Inches(0.12), Inches(2.02), Inches(2.84), Inches(0.90), size=17, color=NAVY, bold=True)
    txt(s, desc,  x+Inches(0.12), Inches(2.97), Inches(2.84), Inches(1.80), size=13, color=GRAY)
    if i < 3:
        txt(s, "→", Inches(3.33 + i * 3.25), Inches(3.35), Inches(0.45), Inches(0.55),
            size=26, color=SEG, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Automated segmentation removes subjectivity and scales to large datasets",
    Inches(0.30), Inches(6.45), Inches(12.73), Inches(0.55),
    size=14, color=SEG, italic=True, align=PP_ALIGN.CENTER)
notes(s, """Segmentation is different from reconstruction — instead of recovering a continuous quantity (like attenuation or magnetisation), we want a discrete label map.
It's still an inverse problem: from the observed image, we want to recover the underlying structure map. And it's ill-posed: multiple acquisition conditions can produce images that look similar but have different underlying anatomy.
Two additional challenges specific to segmentation:
  — Class imbalance: background pixels typically vastly outnumber foreground pixels (e.g. retinal vessels are ~5% of pixels in DRIVE). Standard cross-entropy is dominated by background.
  — Expert annotation cost: every training image requires a doctor to manually draw structure boundaries. This is expensive and slow — DRIVE has only 20 training images.
Automated segmentation matters clinically because it removes inter-observer variability (different radiologists draw different boundaries), speeds up workflows, and enables large-scale population studies.""")

# ── 31. TWO APPROACHES ───────────────────────────────────────────────────────
s = content_slide(SEG, "Two Segmentation Paradigms")
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(5.85), SEG_L)
box(s, Inches(0.30), Inches(1.15), Inches(5.95), Inches(0.10), SEG)
txt(s, "Data-Driven: R2U-Net",
    Inches(0.48), Inches(1.30), Inches(5.65), Inches(0.52), size=19, color=SEG, bold=True)
for i, p in enumerate(["Recurrent Residual U-Net  (Alom et al. 2019)",
                        "Encoder–decoder with skip connections",
                        "Recurrent + residual conv blocks",
                        "Learns pixel-wise mapping from annotated pairs",
                        "Evaluated: retinal vessels, skin & lung lesions",
                        "Outperforms U-Net, SegNet, ResU-Net",
                        "Requires large labelled training set"]):
    txt(s, "•  " + p, Inches(0.48), Inches(1.95 + i * 0.60), Inches(5.65), Inches(0.55), size=13, color=GRAY)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(5.85), MRI_L)
box(s, Inches(7.08), Inches(1.15), Inches(6.05), Inches(0.10), MRI)
txt(s, "Atlas-Based: Multi-Atlas",
    Inches(7.25), Inches(1.30), Inches(5.75), Inches(0.52), size=19, color=MRI, bold=True)
for i, p in enumerate(["Kim et al. 2019  (cardiac PET, rat model)",
                        "Register source image to small atlas library",
                        "Best atlas selected via Pearson correlation",
                        "Transfer pre-defined ROIs to source space",
                        "No large training set required",
                        "r² = 0.92 vs 0.88 for single-atlas method",
                        "Strongly dependent on registration quality"]):
    txt(s, "•  " + p, Inches(7.25), Inches(1.95 + i * 0.60), Inches(5.75), Inches(0.55), size=13, color=GRAY)
txt(s, "vs", Inches(6.20), Inches(3.60), Inches(0.93), Inches(0.55),
    size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
notes(s, """Paper 1 — Alom et al. 2019, R2U-Net:
This is a deep learning method built on the U-Net architecture. They replace standard convolutional blocks with recurrent residual blocks — recurrent layers let each block refine its feature maps over multiple time steps, while residual connections help gradients flow during training. The result is better feature propagation with similar parameter count to standard U-Net. Evaluated on DRIVE (retinal vessels), STARE, and lung/skin lesion datasets. Reported improvements in Dice, sensitivity, specificity vs U-Net baseline.

Paper 2 — Kim et al. 2019, Multi-atlas cardiac PET:
This is a classical processing-based approach. They manually segment the left ventricular myocardium in a small set of PET scans to create an atlas library. For a new scan: register the source to each atlas using affine + nonlinear registration, select the best-matching atlas via Pearson correlation, and transfer the pre-defined ROIs. Evaluated with leave-one-out cross-validation on 12 rats. r²=0.92 for perfusion estimates — better than single-atlas.

The two papers represent fundamentally different philosophies: one learns from data, one uses expert-encoded knowledge.""")

# ── 32. COMPARISON TABLE ─────────────────────────────────────────────────────
s = content_slide(SEG, "Critical Comparison: R2U-Net vs Multi-Atlas")
col_w = [Inches(3.55), Inches(4.50), Inches(4.50)]
col_x = [Inches(0.32), Inches(3.87), Inches(8.37)]
for j, (cw, cx, cc, ch) in enumerate(zip(col_w, col_x,
                                          [NAVY, SEG, MRI],
                                          ["Dimension", "R2U-Net  (DL)", "Multi-Atlas"])):
    box(s, cx, Inches(1.15), cw - Inches(0.05), Inches(0.52), cc)
    txt(s, ch, cx+Inches(0.07), Inches(1.18), cw-Inches(0.12), Inches(0.46),
        size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (prop, dl_v, atlas_v) in enumerate([
    ("Data needed",       "Large labelled dataset",        "Small atlas library only"),
    ("Generalisation",    "Broad  (learned features)",     "Narrow  (atlas domain-limited)"),
    ("Interpretability",  "Low  (black-box CNN)",          "High  (transparent pipeline)"),
    ("Inference speed",   "Fast  (single forward pass)",   "Slow  (registration per image)"),
    ("Training cost",     "High  (GPU + many epochs)",     "None  (no training needed)"),
    ("Boundary quality",  "Can be blurry (Dice/CE loss)",  "Registration-quality limited"),
    ("Best suited for",   "Common tasks + large datasets", "Specialised, small-data settings"),
]):
    y = Inches(1.72 + i * 0.72)
    bg = LGRAY if i % 2 == 0 else WHITE
    for j, (val, cw, cx) in enumerate(zip([prop, dl_v, atlas_v], col_w, col_x)):
        box(s, cx, y, cw - Inches(0.05), Inches(0.69), bg)
        txt(s, val, cx+Inches(0.09), y+Inches(0.06), cw-Inches(0.18), Inches(0.60),
            size=12, color=NAVY if j == 0 else GRAY, bold=(j == 0))
notes(s, """Walking through the comparison table:
Data: R2U-Net needs thousands of annotated images — expensive to obtain. Atlas method needs only a small hand-labelled atlas library.
Generalisation: R2U-Net learns texture, shape, and context from data — it can generalise to varied cases if training is diverse. Atlas method is limited to the domain of its atlases — works poorly if the new image looks different from the library.
Interpretability: R2U-Net is a black box — you can't easily explain why it predicts a particular boundary. Atlas method is transparent: you can trace every decision (which atlas was selected, which registration was applied, how the ROI was transferred).
Inference speed: R2U-Net runs in milliseconds once trained. Atlas method must register the new image to every atlas at inference time — registration is expensive, can take minutes per image.
Boundary quality: interesting tension. R2U-Net's Dice loss optimises region overlap, not boundary sharpness — boundaries tend to be slightly blurred. Atlas method's boundary quality depends entirely on registration accuracy — if registration fails, boundaries are wrong in a systematic way.
My conclusion: R2U-Net is better for common, large-dataset segmentation tasks (retinal vessels, liver, etc.). Multi-atlas is better for specialised, small-data clinical settings where interpretability and regulatory approval matter (like the cardiac PET study).""")

# ── 33. IMPROVEMENTS ─────────────────────────────────────────────────────────
s = content_slide(SEG, "Proposed Improvements to R2U-Net  (Alom et al. 2019)",
                  "Three complementary extensions  ·  evaluate on same datasets with ablation study")
for i, (color, title, desc) in enumerate([
    (SEG, "Attention-Gated\nSkip Connections",
     "Plain skips propagate irrelevant background\nfeatures into decoder alongside useful ones.\n\n"
     "Attention gates learn soft weights guided by\ndecoder context, suppressing spatially irrelevant\n"
     "activations before concatenation.\n\n"
     "Retains recurrent-residual structure intact.\nMinimal parameter overhead."),
    (MRI, "Hausdorff Distance\nLoss Component",
     "Dice + cross-entropy optimise region overlap\nbut do not penalise boundary displacement.\n\n"
     "Adding a Hausdorff-distance loss term encourages\nthe network to reduce large localised errors\n"
     "near object boundaries.\n\n"
     "Sharpens predicted contours without sacrificing\noverall segmentation accuracy."),
    (CT,  "GAN-Based\nData Augmentation",
     "Data scarcity: DRIVE has only 20 training images.\nStandard geometric transforms lack variety.\n\n"
     "Conditional GAN generates plausible retinal /\nlesion images with paired segmentation masks,\n"
     "introducing anatomical variety beyond transforms.\n\n"
     "Tested with leave-one-out on same splits;\ntrack Dice, sensitivity, AUC + Hausdorff distance."),
]):
    x = Inches(0.30 + i * 4.34)
    w = Inches(4.12)
    box(s, x, Inches(1.15), w, Inches(6.18), LGRAY)
    box(s, x, Inches(1.15), w, Inches(0.10), color)
    txt(s, title, x+Inches(0.12), Inches(1.30), w-Inches(0.25), Inches(0.90), size=16, color=color, bold=True)
    txt(s, desc,  x+Inches(0.12), Inches(2.25), w-Inches(0.25), Inches(4.90), size=12, color=GRAY)
notes(s, """Three improvements I proposed, each targeting a specific weakness in R2U-Net:

1. Attention-gated skip connections (motivated by Schlemper et al. 2019, Attention U-Net):
R2U-Net uses plain skip connections that forward all encoder feature maps to the decoder. In retinal vessel segmentation, where vessels are thin and surrounded by vast background, irrelevant background activations are forwarded equally alongside vessel features. Attention gates compute a soft spatial weight guided by the decoder's context, suppressing background features before they reach the decoder. This keeps the recurrent-residual structure intact and adds only a small number of parameters.

2. Hausdorff distance loss:
Standard Dice and cross-entropy losses optimise for region overlap — they penalise every misclassified pixel equally. But clinically, a large boundary error (e.g. being 10 pixels away from the true vessel edge) is much worse than many small distributed errors. The Hausdorff distance measures the worst-case boundary displacement. Adding it to the loss pushes the network to reduce catastrophic localised boundary misses.

3. GAN-based data augmentation:
DRIVE has 20 training images — far too few for training a robust network. Standard augmentation (rotation, flip, elastic deformation) helps but doesn't add anatomical variety. A conditional GAN trained on the same 20 images can synthesise new plausible retinal images with paired segmentation masks, introducing more varied vessel patterns and textures. This is especially valuable because expert annotation costs prevent collecting more real data.""")

# ── 34. FINAL TAKEAWAYS ──────────────────────────────────────────────────────
s = new_slide()
set_bg(s, NAVY)
box(s, 0, 0, W, Inches(0.08), CT)
txt(s, "Three Key Takeaways",
    Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.72),
    size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(3.8), Inches(0.92), Inches(5.73), Inches(0.05), RGBColor(0x30, 0x60, 0x90))
for i, (color, num, title, desc) in enumerate([
    (CT, "1",
     "Reconstruction quality depends on both data and algorithm",
     "High dose + dense views is required for FBP to work.\n"
     "GD degrades gracefully: still produces useful images at I₀=10² "
     "where FBP collapses to −5.7 dB."),
    (MRI, "2",
     "Iterative methods are robust — but require careful convergence control",
     "GD outperforms FBP in all 9 CT conditions.\n"
     "Subset GD accelerates convergence but diverges without "
     "early stopping or step-size decay."),
    (SEG, "3",
     "Denoising domain and modelling choices matter",
     "k-space vs image-space filtering produce qualitatively different results — "
     "bilateral preserves edges best.\n"
     "For segmentation: DL generalises broadly; atlas methods excel in small-data settings."),
]):
    y = Inches(1.15 + i * 1.93)
    box(s, Inches(0.30), y, Inches(0.78), Inches(0.78), color)
    txt(s, num, Inches(0.30), y, Inches(0.78), Inches(0.78),
        size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(1.18), y, Inches(11.90), Inches(1.80), NAVY2)
    box(s, Inches(1.18), y, Inches(0.07), Inches(1.80), color)
    txt(s, title, Inches(1.33), y + Inches(0.10), Inches(11.65), Inches(0.50),
        size=18, color=color, bold=True)
    txt(s, desc,  Inches(1.33), y + Inches(0.60), Inches(11.65), Inches(1.10),
        size=13, color=RGBColor(0xcc, 0xd3, 0xe0))
box(s, 0, H - Inches(0.05), W, Inches(0.05), CT)
txt(s, "Medical Imaging Coursework  ·  Barbara Koch  ·  "
       "MPhil DIS, University of Cambridge  ·  Lent 2026",
    Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.30),
    size=11, color=RGBColor(0x55, 0x66, 0x77), align=PP_ALIGN.CENTER)
notes(s, """Wrapping up with three messages that cut across all modules.

1. Data quality AND algorithm together determine the result. The most important finding from Module 1: GD doesn't magically create information — but it uses the available noisy data more efficiently than FBP. At I₀=10², GD gives 17.3 dB vs FBP's 7.2 dB at 360 views. That's the difference between a clinically usable image and noise. Algorithm choice matters enormously in low-data regimes.

2. Iterative methods are more robust but need careful handling. GD wins across all 9 conditions, but naive unregularised GD still has semi-convergence: it peaks then slowly deteriorates. Subset GD converges 3× faster but diverges without early stopping. The practical lesson: iterative methods need regularisation (TV, early stopping, step decay) to be clinically deployed safely.

3. Choices throughout the pipeline compound. In MRI: whether you filter in k-space or image space changes what information you lose (global vs local). Within image space: bilateral preserves edges; Gaussian blurs them. In segmentation: the loss function (Dice vs Hausdorff-augmented) changes what kind of errors the model makes. Getting the best result requires understanding your modality, your noise model, and your clinical goal at every step.

Thank you. Happy to take questions on any of the three modules.""")

# ===========================================================================
# SAVE
# ===========================================================================
out = BASE + "/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
